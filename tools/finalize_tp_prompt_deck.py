#!/usr/bin/env python3
from __future__ import annotations

import copy
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
BACKUP = ROOT / "docs/source/TP Prompt Ingénierie.before_codex_merge.bak.pptx"
CODEX = ROOT / "artifacts/codex_cli_skills_performance_fr.pptx"
TARGET = ROOT / "docs/source/TP Prompt Ingénierie.pptx"

KEEP_BACKUP = [1, 4, 5, 8, 10]
KEEP_CODEX = [2, 3, 5]

BG = RGBColor(247, 243, 235)
INK = RGBColor(33, 36, 39)
MUTED = RGBColor(91, 98, 104)
ACCENT = RGBColor(184, 87, 36)
ACCENT_2 = RGBColor(45, 92, 122)
PANEL = RGBColor(255, 251, 245)
LINE = RGBColor(220, 206, 191)


def clone_slide(src_slide, dst_prs):
    dst_slide = dst_prs.slides.add_slide(dst_prs.slide_layouts[6])
    for shape in src_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        dst_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")  # pyright: ignore
    return dst_slide


def iter_text_shapes(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_band(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def style_shape(shape, size: int, color: RGBColor, bold: bool = False) -> None:
    tf = shape.text_frame
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Aptos"
        if not paragraph.runs and paragraph.text:
            run = paragraph.add_run()
            run.text = paragraph.text
            paragraph.text = ""
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Aptos"


def style_slide(slide, idx: int) -> None:
    set_bg(slide)
    add_top_band(slide)
    text_shapes = list(iter_text_shapes(slide))
    if not text_shapes:
        return
    style_shape(text_shapes[0], 28 if idx == 1 else 24, ACCENT_2, True)
    for shape in text_shapes[1:]:
        style_shape(shape, 16 if idx == 1 else 18, MUTED if idx == 1 else INK, False)


def tighten_codex(slide, idx: int) -> None:
    texts = list(iter_text_shapes(slide))
    if idx == 6 and len(texts) >= 5:
        texts[3].text_frame.text = "UI AGILab vaste\npages, apps, runs"
        style_shape(texts[3], 22, INK, False)
        texts[4].text_frame.text = "Le skill evite de chercher partout."
        style_shape(texts[4], 14, MUTED, False)
    elif idx == 7 and len(texts) >= 4:
        texts[3].text_frame.text = "metadonnees\nSKILL.md\nscripts utiles"
        style_shape(texts[3], 22, INK, False)
    elif idx == 8 and len(texts) >= 4:
        texts[3].text_frame.text = "AGENTS.md\nSkill\nPlugin\nMCP"
        style_shape(texts[3], 22, INK, False)


def add_takeaways(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(11.5), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "10. Takeaways"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2
    r.font.name = "Aptos"
    box = slide.shapes.add_textbox(Inches(0.72), Inches(1.3), Inches(11.2), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Trois idees a retenir"
    r.font.size = Pt(13)
    r.font.color.rgb = MUTED
    r.font.name = "Aptos"
    cards = [
        (0.85, "1", "Le skill n'ajoute pas plus de modele.\nIl ajoute un meilleur workflow."),
        (4.45, "2", "Dans AGILab, il rappelle comment lancer,\nverifier et ne pas casser le core."),
        (8.05, "3", "On les met a jour volontairement,\npour garder un repo stable et tracable."),
    ]
    for left, head, body in cards:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.2), Inches(2.95), Inches(3.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = LINE
        box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.42), Inches(0.5), Inches(0.4))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = head
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        r.font.name = "Aptos"
        box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(3.0), Inches(2.4), Inches(1.7))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = body
        r.font.size = Pt(18)
        r.font.color.rgb = INK
        r.font.name = "Aptos"


def add_compare_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide)

    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(11.5), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "9. Code Companion V3 vs Codex CLI"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2
    r.font.name = "Aptos"

    box = slide.shapes.add_textbox(Inches(0.72), Inches(1.3), Inches(11.2), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Deux outils complementaires dans AGILab"
    r.font.size = Pt(13)
    r.font.color.rgb = MUTED
    r.font.name = "Aptos"

    cards = [
        (
            0.85,
            "Code Companion V3",
            "Integre et securise\nAzure entreprise\nPyCharm + VS Code\nIdeal pour l'iteration locale",
        ),
        (
            6.75,
            "Codex CLI",
            "Agent en terminal\nskills + AGENTS.md\nworkflows repo\nIdeal pour taches longues",
        ),
    ]
    for left, head, body in cards:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(2.15),
            Inches(5.0),
            Inches(3.2),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = LINE

        box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(2.45), Inches(4.2), Inches(0.45))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = head
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        r.font.name = "Aptos"

        box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(3.05), Inches(4.2), Inches(1.8))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = body
        r.font.size = Pt(18)
        r.font.color.rgb = INK
        r.font.name = "Aptos"

    box = slide.shapes.add_textbox(Inches(1.1), Inches(5.85), Inches(10.6), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Lecture simple : IDE securise d'un cote, agent repo de l'autre."
    r.font.size = Pt(16)
    r.font.color.rgb = MUTED
    r.font.name = "Aptos"


def main() -> int:
    backup = Presentation(str(BACKUP))
    codex = Presentation(str(CODEX))
    prs = Presentation()

    while len(prs.slides):
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])  # pyright: ignore

    for idx in KEEP_BACKUP:
        clone_slide(backup.slides[idx - 1], prs)
    for idx in KEEP_CODEX:
        clone_slide(codex.slides[idx - 1], prs)

    for idx, slide in enumerate(prs.slides, start=1):
        style_slide(slide, idx)
        if idx >= 6:
            tighten_codex(slide, idx)

    add_compare_slide(prs)
    add_takeaways(prs)

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir) / TARGET.name
        prs.save(str(tmp))
        TARGET.write_bytes(tmp.read_bytes())

    print(TARGET)
    print("slides", len(prs.slides))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
