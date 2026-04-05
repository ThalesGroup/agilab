#!/usr/bin/env python3
from __future__ import annotations

import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TARGET = ROOT / "docs/source/TP Prompt Ingénierie.pptx"

BG = RGBColor(247, 243, 235)
INK = RGBColor(33, 36, 39)
MUTED = RGBColor(91, 98, 104)
ACCENT = RGBColor(184, 87, 36)
ACCENT_2 = RGBColor(45, 92, 122)
PANEL = RGBColor(255, 251, 245)
LINE = RGBColor(220, 206, 191)


def _iter_text_shapes(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape


def _set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _has_top_band(slide) -> bool:
    for shape in slide.shapes:
        if getattr(shape, "shape_type", None) != 1:
            continue
        if getattr(shape, "top", 1) == 0 and getattr(shape, "height", 0) < 350000:
            return True
    return False


def _add_top_band(slide) -> None:
    if _has_top_band(slide):
        return
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.22)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def _style_runs(shape, *, size: int, color: RGBColor, bold: bool = False, name: str = "Aptos") -> None:
    tf = shape.text_frame
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = name
        if not paragraph.runs and paragraph.text:
            run = paragraph.add_run()
            run.text = paragraph.text
            paragraph.text = ""
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = name


def _style_standard_slide(slide, slide_number: int) -> None:
    _set_bg(slide)
    _add_top_band(slide)

    text_shapes = list(_iter_text_shapes(slide))
    if not text_shapes:
        return

    title = text_shapes[0]
    _style_runs(title, size=24 if slide_number > 1 else 28, color=ACCENT_2, bold=True)
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT

    for shape in text_shapes[1:]:
        _style_runs(shape, size=18 if slide_number > 1 else 16, color=INK, bold=False)

    # Title slide special case.
    if slide_number == 1 and len(text_shapes) > 1:
        _style_runs(text_shapes[1], size=16, color=MUTED, bold=False)


def _find_by_text(slide, needle: str):
    for shape in _iter_text_shapes(slide):
        full = " ".join(p.text for p in shape.text_frame.paragraphs)
        if needle in full:
            return shape
    return None


def _tighten_codex_slides(prs: Presentation) -> None:
    replacements = {
        7: ("UI AGILab vaste\npages, apps, runs", "Le skill evite de chercher partout."),
        8: ("metadonnees\nSKILL.md\nscripts utiles", None),
        9: ("AGENTS.md\nSkill\nPlugin\nMCP", None),
        10: ("scripts\nreferences\nassets", "Fait / cadre / fournit"),
        11: ("pas d'auto-rewrite\nstabilite\ntracabilite", None),
        12: ("Codex devient plus fiable\ncar il apprend\ncomment travailler ici.", "Les skills restent des artefacts humains."),
    }

    for slide_idx, (main_text, note_text) in replacements.items():
        slide = prs.slides[slide_idx - 1]
        panel_text = None
        candidates = [s for s in _iter_text_shapes(slide)]
        if slide_idx in {7, 8, 9, 10, 11} and len(candidates) >= 8:
            panel_text = candidates[7]
            tf = panel_text.text_frame
            tf.clear()
            for i, line in enumerate(main_text.split("\n")):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = INK
                p.font.name = "Aptos"
                p.space_after = Pt(10)
                p.bullet = False
            if note_text and len(candidates) >= 9:
                note = candidates[-1]
                note.text_frame.clear()
                p = note.text_frame.paragraphs[0]
                p.text = note_text
                p.font.size = Pt(14)
                p.font.color.rgb = MUTED
                p.font.name = "Aptos"
        elif slide_idx == 12 and len(candidates) >= 8:
            main = candidates[6]
            sub = candidates[7]
            main.text_frame.clear()
            p = main.text_frame.paragraphs[0]
            p.text = main_text
            p.font.size = Pt(24)
            p.font.color.rgb = INK
            p.font.bold = True
            p.font.name = "Aptos"
            sub.text_frame.clear()
            p = sub.text_frame.paragraphs[0]
            p.text = note_text or ""
            p.font.size = Pt(16)
            p.font.color.rgb = MUTED
            p.font.name = "Aptos"


def _add_takeaways(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_top_band(slide)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.7), Inches(11.5), Inches(0.6))
    p = title.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "10. Takeaways"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = ACCENT_2
    r.font.name = "Aptos"

    subtitle = slide.shapes.add_textbox(Inches(0.72), Inches(1.35), Inches(11.0), Inches(0.4))
    p = subtitle.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "Trois messages a retenir pour AGILab"
    r.font.size = Pt(13)
    r.font.color.rgb = MUTED
    r.font.name = "Aptos"

    for left, top, head, body in [
        (0.9, 2.2, "1", "Un skill n'ajoute pas plus de modele.\nIl ajoute un meilleur workflow."),
        (4.45, 2.2, "2", "Dans AGILab, il rappelle comment lancer,\nverifier et ne pas casser le core."),
        (8.0, 2.2, "3", "On les met a jour volontairement,\npour garder un repo stable et tracable."),
    ]:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(2.95),
            Inches(3.0),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PANEL
        shape.line.color.rgb = LINE

        box = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.2), Inches(2.4), Inches(0.5))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = head
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        r.font.name = "Aptos"

        box = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.8), Inches(2.45), Inches(1.8))
        p = box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = body
        r.font.size = Pt(18)
        r.font.color.rgb = INK
        r.font.name = "Aptos"


def main() -> int:
    prs = Presentation(str(TARGET))

    for idx, slide in enumerate(prs.slides, start=1):
        _style_standard_slide(slide, idx)

    _tighten_codex_slides(prs)
    _add_takeaways(prs)

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir) / TARGET.name
        prs.save(str(tmp))
        TARGET.write_bytes(tmp.read_bytes())

    print(TARGET)
    print("slides", len(prs.slides))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
