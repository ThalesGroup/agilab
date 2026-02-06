#!/usr/bin/env python3
"""
Lightweight PPTX alignment helper.

Purpose
-------
Some decks are built with a custom "header text box" near the top of each slide,
plus a large placeholder used for the main content. When shapes are nudged around
manually, slides become visually inconsistent (different margins, slight negative
coordinates, overlapping header/content).

This script applies conservative, repeatable alignment:
  1) Detect a "header" textbox near the top of each slide.
  2) Move headers to a canonical (median) left/top position across the deck.
  3) Snap the main TITLE placeholder left margin to 0 (removes small negatives).
  4) For slides where the content box is already close to the header, enforce a
     consistent vertical gap between header and content (median gap).

It is intentionally heuristic and avoids resizing shapes to reduce text reflow.
"""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path
from statistics import median
from typing import Iterable

from pptx import Presentation


@dataclass
class _SlideTargets:
    slide_index: int
    header_shape: object | None
    body_placeholder: object | None
    header_left: int | None
    header_top: int | None
    header_height: int | None
    body_left: int | None
    body_top: int | None


def _iter_slides(prs: Presentation) -> Iterable[tuple[int, object]]:
    for idx, slide in enumerate(prs.slides, start=1):
        yield idx, slide


def _shape_text(shape: object) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    text_frame = shape.text_frame
    return (text_frame.text or "").strip()


def _looks_like_header_text(text: str) -> bool:
    if not text:
        return False
    stripped = text.lstrip()
    # Heuristic: body bullets often start with a bullet symbol.
    if stripped.startswith("•"):
        return False
    return True


def _find_header_shape(slide: object, *, max_top_emu: int) -> object | None:
    candidates: list[object] = []
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _shape_text(shape)
        if not _looks_like_header_text(text):
            continue
        if getattr(shape, "top", 0) > max_top_emu:
            continue
        candidates.append(shape)
    if not candidates:
        return None
    # Pick the top-most (then left-most) candidate.
    candidates.sort(key=lambda s: (getattr(s, "top", 0), getattr(s, "left", 0)))
    return candidates[0]


def _find_title_placeholder(slide: object) -> object | None:
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            # 1 == TITLE in pptx.enum.shapes.PP_PLACEHOLDER
            if shape.placeholder_format.type == 1:
                return shape
        except Exception:
            continue
    return None


def _snap_small_negative(value: int, *, epsilon: int) -> int:
    return 0 if -epsilon <= value < 0 else value


def main() -> int:
    parser = argparse.ArgumentParser(description="Conservatively align PPTX slides.")
    parser.add_argument("--in", dest="input_path", required=True, help="Input .pptx")
    parser.add_argument("--out", dest="output_path", required=True, help="Output .pptx")
    parser.add_argument(
        "--header-max-top-emu",
        type=int,
        default=600_000,
        help="Only consider textboxes above this top value as headers (EMU).",
    )
    parser.add_argument(
        "--body-gap-threshold-emu",
        type=int,
        default=300_000,
        help="Only enforce header/body gap when the original gap is <= this (EMU).",
    )
    parser.add_argument(
        "--snap-negative-epsilon-emu",
        type=int,
        default=10_000,
        help="Snap slightly negative coordinates within this epsilon to 0 (EMU).",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Print detected medians and planned edits, but do not write output.",
    )
    args = parser.parse_args()

    input_path = Path(args.input_path).expanduser()
    output_path = Path(args.output_path).expanduser()

    prs = Presentation(str(input_path))

    slide_targets: list[_SlideTargets] = []
    header_lefts: list[int] = []
    header_tops: list[int] = []
    positive_gaps: list[int] = []

    # Collect per-slide targets + medians.
    for slide_index, slide in _iter_slides(prs):
        header = _find_header_shape(slide, max_top_emu=args.header_max_top_emu)
        body = _find_title_placeholder(slide)

        header_left = getattr(header, "left", None) if header else None
        header_top = getattr(header, "top", None) if header else None
        header_height = getattr(header, "height", None) if header else None
        body_left = getattr(body, "left", None) if body else None
        body_top = getattr(body, "top", None) if body else None

        if slide_index != 1 and header_left is not None and header_top is not None:
            header_lefts.append(int(header_left))
            header_tops.append(int(header_top))

        if (
            slide_index != 1
            and header is not None
            and body is not None
            and header_top is not None
            and header_height is not None
            and body_top is not None
        ):
            gap = int(body_top) - (int(header_top) + int(header_height))
            if gap > 0:
                positive_gaps.append(gap)

        slide_targets.append(
            _SlideTargets(
                slide_index=slide_index,
                header_shape=header,
                body_placeholder=body,
                header_left=int(header_left) if header_left is not None else None,
                header_top=int(header_top) if header_top is not None else None,
                header_height=int(header_height) if header_height is not None else None,
                body_left=int(body_left) if body_left is not None else None,
                body_top=int(body_top) if body_top is not None else None,
            )
        )

    if not header_lefts or not header_tops:
        raise RuntimeError("Could not detect any header shapes to align.")

    canonical_header_left = int(median(header_lefts))
    canonical_header_top = int(median(header_tops))
    canonical_gap = int(median(positive_gaps)) if positive_gaps else 150_000

    if args.dry_run:
        print("canonical_header_left:", canonical_header_left)
        print("canonical_header_top :", canonical_header_top)
        print("canonical_gap        :", canonical_gap)

    # Apply edits.
    for target in slide_targets:
        if target.slide_index == 1:
            continue

        header = target.header_shape
        if header is not None:
            header.left = canonical_header_left
            header.top = canonical_header_top
            header.left = _snap_small_negative(
                int(header.left), epsilon=args.snap_negative_epsilon_emu
            )
            header.top = _snap_small_negative(
                int(header.top), epsilon=args.snap_negative_epsilon_emu
            )

        body = target.body_placeholder
        if body is not None:
            # Normalise small negative left coordinates.
            body.left = _snap_small_negative(
                int(body.left), epsilon=args.snap_negative_epsilon_emu
            )
            # Most decks expect the main content box to start at the left edge.
            body.left = 0

            if header is not None and target.header_top is not None and target.body_top is not None:
                # Use original gap to decide whether this slide is "close enough" to normalise.
                original_gap = target.body_top - (target.header_top + (target.header_height or 0))
                if original_gap <= args.body_gap_threshold_emu:
                    desired_top = int(header.top) + int(header.height) + canonical_gap
                    body.top = max(desired_top, 0)

    if args.dry_run:
        return 0

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

