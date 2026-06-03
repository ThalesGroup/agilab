#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path
import shutil

from pptx import Presentation

from merge_codex_skills_into_tp_prompt import (
    ADD_SOURCE,
    BACKUP,
    SOURCE,
    _clone_slide_content,
    _remove_slide,
    _slide_text,
)


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs/source/TP Prompt Ingénierie.rebuilt.pptx"
TARGET = ROOT / "docs/source/TP Prompt Ingénierie.pptx"
REMOVE_TARGET = [2, 3, 7, 9, 11]


def build(out_path: Path) -> None:
    shutil.copy2(BACKUP, out_path)
    target = Presentation(str(out_path))
    source = Presentation(str(SOURCE))

    for one_based in sorted(REMOVE_TARGET, reverse=True):
        _remove_slide(target, one_based - 1)
    target.save(str(out_path))

    target = Presentation(str(out_path))
    blank_layout = target.slide_layouts[6]
    for one_based in ADD_SOURCE:
        src_slide = source.slides[one_based - 1]
        dst_slide = target.slides.add_slide(blank_layout)
        _clone_slide_content(src_slide, dst_slide)
    target.save(str(out_path))


def main() -> int:
    build(OUT)
    prs = Presentation(str(OUT))
    print("out:", OUT)
    print("slides:", len(prs.slides))
    for i in range(7, 13):
        print(i, _slide_text(prs.slides[i - 1])[:180])
    shutil.copy2(OUT, TARGET)
    print("replaced:", TARGET)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
