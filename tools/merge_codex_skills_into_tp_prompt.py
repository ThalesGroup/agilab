#!/usr/bin/env python3
from __future__ import annotations

import copy
from pathlib import Path
import tempfile

from pptx import Presentation
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
TARGET = ROOT / "docs/source/TP Prompt Ingénierie.pptx"
SOURCE = ROOT / "artifacts/codex_cli_skills_performance_fr.pptx"
BACKUP = ROOT / "docs/source/TP Prompt Ingénierie.before_codex_merge.bak.pptx"

# 1-based slide numbers from the target deck to remove because they add little value
# relative to the new Codex material.
REMOVE_TARGET = [2, 3, 7, 9, 11]

# 1-based slide numbers from the generated Codex deck to append.
ADD_SOURCE = [2, 3, 5, 6, 8, 10]


def _slide_text(slide) -> str:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = " ".join(
                p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
            ).strip()
            if text:
                texts.append(text)
    return " | ".join(texts)


def _remove_slide(prs: Presentation, idx: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # pyright: ignore[reportAttributeAccessIssue]
    slides = list(slide_id_list)
    slide_id = slides[idx]
    r_id = slide_id.rId
    prs.part.drop_rel(r_id)
    slide_id_list.remove(slide_id)


def _copy_background(src_slide, dst_slide) -> None:
    src_fill = src_slide.background.fill
    if src_fill.type == 1:
        dst_fill = dst_slide.background.fill
        dst_fill.solid()
        try:
            dst_fill.fore_color.rgb = RGBColor.from_string(str(src_fill.fore_color.rgb))
        except Exception:
            pass


def _clone_slide_content(src_slide, dst_slide) -> None:
    _copy_background(src_slide, dst_slide)
    for shape in src_slide.shapes:
        new_el = copy.deepcopy(shape.element)
        dst_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")  # pyright: ignore[reportAttributeAccessIssue]


def main() -> int:
    if not TARGET.exists():
        raise FileNotFoundError(TARGET)
    if not SOURCE.exists():
        raise FileNotFoundError(SOURCE)

    if not BACKUP.exists():
        BACKUP.write_bytes(TARGET.read_bytes())

    # Start from the preserved clean backup each time so reruns are deterministic,
    # but work on temporary files to avoid rewriting the same zip path in place.
    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)
        phase1 = tmpdir / "phase1_pruned.pptx"
        phase2 = tmpdir / "phase2_merged.pptx"
        phase1.write_bytes(BACKUP.read_bytes())

        target = Presentation(str(phase1))
        source = Presentation(str(SOURCE))

        removed = []
        for one_based in sorted(REMOVE_TARGET, reverse=True):
            removed.append((one_based, _slide_text(target.slides[one_based - 1])[:140]))
            _remove_slide(target, one_based - 1)

        target.save(str(phase1))

        # Reopen after pruning so new slide part names are generated cleanly.
        target = Presentation(str(phase1))
        blank_layout = target.slide_layouts[6]
        added = []
        for one_based in ADD_SOURCE:
            src_slide = source.slides[one_based - 1]
            dst_slide = target.slides.add_slide(blank_layout)
            _clone_slide_content(src_slide, dst_slide)
            added.append((one_based, _slide_text(src_slide)[:140]))

        target.save(str(phase2))
        TARGET.write_bytes(phase2.read_bytes())

    print("backup:", BACKUP)
    print("target:", TARGET)
    print("removed:")
    for idx, text in removed:
        print(f"  - {idx}: {text}")
    print("added:")
    for idx, text in added:
        print(f"  - {idx}: {text}")
    print("final_slide_count:", len(target.slides))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
