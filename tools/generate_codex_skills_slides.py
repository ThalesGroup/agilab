#!/usr/bin/env python3
from __future__ import annotations

import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ROOT = Path(__file__).resolve().parents[1]
DIAGRAM_DIR = ROOT / "docs/source/diagrams/codex_skills"
OUTPUT = ROOT / "artifacts/codex_cli_skills_performance_fr.pptx"

BG = RGBColor(247, 243, 235)
INK = RGBColor(33, 36, 39)
MUTED = RGBColor(91, 98, 104)
ACCENT = RGBColor(184, 87, 36)
ACCENT_2 = RGBColor(45, 92, 122)
PANEL = RGBColor(255, 251, 245)
LINE = RGBColor(220, 206, 191)


def _write(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def _render_png(svg_path: Path) -> Path:
    png_path = svg_path.with_suffix(".png")
    subprocess.run(
        ["/opt/homebrew/bin/rsvg-convert", str(svg_path), "-w", "2400", "-o", str(png_path)],
        check=True,
    )
    return png_path


def _build_diagrams() -> dict[str, Path]:
    diagrams: dict[str, str] = {
        "story": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="72" y="118" width="262" height="182" rx="30" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="469" y="118" width="262" height="182" rx="30" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="866" y="118" width="262" height="182" rx="30" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <circle cx="118" cy="92" r="26" fill="#b85724"/><text x="118" y="101" font-size="24" text-anchor="middle" fill="#fffaf5" font-family="Aptos" font-weight="700">1</text>
  <circle cx="515" cy="92" r="26" fill="#b85724"/><text x="515" y="101" font-size="24" text-anchor="middle" fill="#fffaf5" font-family="Aptos" font-weight="700">2</text>
  <circle cx="912" cy="92" r="26" fill="#b85724"/><text x="912" y="101" font-size="24" text-anchor="middle" fill="#fffaf5" font-family="Aptos" font-weight="700">3</text>
  <text x="203" y="185" font-size="32" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">AGILab</text>
  <text x="203" y="228" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">pages, apps, runs</text>
  <text x="203" y="260" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">beaucoup de chemins possibles</text>
  <text x="600" y="185" font-size="32" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Codex</text>
  <text x="600" y="228" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">sélectionne le bon workflow</text>
  <text x="600" y="260" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">sans relire tout le repo</text>
  <text x="997" y="185" font-size="32" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Resultat</text>
  <text x="997" y="228" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">patch plus fiable</text>
  <text x="997" y="260" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">tests et conventions respectes</text>
  <path d="M334 210 L469 210" stroke="#b85724" stroke-width="8" fill="none" stroke-linecap="round"/>
  <path d="M731 210 L866 210" stroke="#b85724" stroke-width="8" fill="none" stroke-linecap="round"/>
  <polygon points="455,198 480,210 455,222" fill="#b85724"/>
  <polygon points="852,198 877,210 852,222" fill="#b85724"/>
  <rect x="150" y="392" width="900" height="156" rx="30" fill="#efe2d1" stroke="none"/>
  <text x="600" y="458" font-size="36" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">Message</text>
  <text x="600" y="505" font-size="25" text-anchor="middle" fill="#212427" font-family="Aptos">Les skills reduisent l'espace de recherche avant d'agir.</text>
</svg>""",
        "flow": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="62" y="142" width="228" height="146" rx="26" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="334" y="142" width="228" height="146" rx="26" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="606" y="142" width="228" height="146" rx="26" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="878" y="142" width="228" height="146" rx="26" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <text x="176" y="188" font-size="30" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">1</text>
  <text x="448" y="188" font-size="30" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">2</text>
  <text x="720" y="188" font-size="30" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">3</text>
  <text x="992" y="188" font-size="30" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">4</text>
  <text x="176" y="230" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">Metadonnees</text>
  <text x="448" y="230" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">SKILL.md</text>
  <text x="720" y="230" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">references</text>
  <text x="720" y="258" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">et scripts utiles</text>
  <text x="992" y="230" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">Execution</text>
  <path d="M290 214 L334 214" stroke="#b85724" stroke-width="8" fill="none" stroke-linecap="round"/>
  <path d="M562 214 L606 214" stroke="#b85724" stroke-width="8" fill="none" stroke-linecap="round"/>
  <path d="M834 214 L878 214" stroke="#b85724" stroke-width="8" fill="none" stroke-linecap="round"/>
  <polygon points="324,202 349,214 324,226" fill="#b85724"/>
  <polygon points="596,202 621,214 596,226" fill="#b85724"/>
  <polygon points="868,202 893,214 868,226" fill="#b85724"/>
  <rect x="152" y="406" width="896" height="150" rx="30" fill="#efe2d1" stroke="none"/>
  <text x="600" y="470" font-size="42" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">Progressive disclosure</text>
  <text x="600" y="518" font-size="25" text-anchor="middle" fill="#212427" font-family="Aptos">On ne charge que ce qui est utile, au moment utile.</text>
</svg>""",
        "agilab_case": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="78" y="112" width="324" height="212" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="438" y="112" width="324" height="212" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="798" y="112" width="324" height="212" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <text x="240" y="176" font-size="29" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Demande</text>
  <text x="240" y="222" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">corriger une page</text>
  <text x="240" y="254" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">ou un formulaire</text>
  <text x="240" y="286" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">sans casser AGILab</text>
  <text x="600" y="176" font-size="29" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Guidage</text>
  <text x="600" y="222" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">AGENTS.md + skills</text>
  <text x="600" y="254" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">+ runbook du repo</text>
  <text x="600" y="286" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">uv, logs, tests, guardrails</text>
  <text x="960" y="176" font-size="29" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Sortie</text>
  <text x="960" y="222" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">patch local</text>
  <text x="960" y="254" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">validation ciblee</text>
  <text x="960" y="286" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">sans toucher le core</text>
  <path d="M402 218 L438 218" stroke="#b85724" stroke-width="8" fill="none" stroke-linecap="round"/>
  <path d="M762 218 L798 218" stroke="#b85724" stroke-width="8" fill="none" stroke-linecap="round"/>
  <polygon points="426,206 451,218 426,230" fill="#b85724"/>
  <polygon points="786,206 811,218 786,230" fill="#b85724"/>
  <rect x="164" y="396" width="872" height="166" rx="30" fill="#efe2d1" stroke="none"/>
  <text x="600" y="463" font-size="34" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">Exemple AGILab</text>
  <text x="600" y="510" font-size="24" text-anchor="middle" fill="#212427" font-family="Aptos">Le skill ne code pas a la place du modele.</text>
  <text x="600" y="545" font-size="24" text-anchor="middle" fill="#212427" font-family="Aptos">Il lui rappelle comment bien travailler dans ce repo.</text>
</svg>""",
        "layers": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="120" y="92" width="960" height="112" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="186" y="230" width="828" height="112" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="252" y="368" width="696" height="112" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="318" y="506" width="564" height="112" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <text x="600" y="160" font-size="34" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">AGENTS.md</text>
  <text x="600" y="197" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">regles durables du repo</text>
  <text x="600" y="298" font-size="34" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Skill</text>
  <text x="600" y="335" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">workflow réutilisable</text>
  <text x="600" y="436" font-size="34" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Plugin</text>
  <text x="600" y="473" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">distribution et partage</text>
  <text x="600" y="574" font-size="34" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">MCP</text>
  <text x="600" y="611" font-size="23" text-anchor="middle" fill="#212427" font-family="Aptos">outils et donnees externes</text>
</svg>""",
        "folders": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="84" y="118" width="304" height="304" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="448" y="118" width="304" height="304" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="812" y="118" width="304" height="304" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <text x="236" y="192" font-size="36" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">scripts/</text>
  <text x="236" y="242" font-size="25" text-anchor="middle" fill="#212427" font-family="Aptos">fait</text>
  <text x="236" y="294" font-size="22" text-anchor="middle" fill="#5b6268" font-family="Aptos">CLI deterministe</text>
  <text x="236" y="328" font-size="22" text-anchor="middle" fill="#5b6268" font-family="Aptos">run_eval.py</text>
  <text x="600" y="192" font-size="36" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">references/</text>
  <text x="600" y="242" font-size="25" text-anchor="middle" fill="#212427" font-family="Aptos">cadre</text>
  <text x="600" y="294" font-size="22" text-anchor="middle" fill="#5b6268" font-family="Aptos">protocoles</text>
  <text x="600" y="328" font-size="22" text-anchor="middle" fill="#5b6268" font-family="Aptos">checklists</text>
  <text x="964" y="192" font-size="36" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">assets/</text>
  <text x="964" y="242" font-size="25" text-anchor="middle" fill="#212427" font-family="Aptos">fournit</text>
  <text x="964" y="294" font-size="22" text-anchor="middle" fill="#5b6268" font-family="Aptos">templates</text>
  <text x="964" y="328" font-size="22" text-anchor="middle" fill="#5b6268" font-family="Aptos">inputs exemples</text>
  <rect x="150" y="462" width="900" height="132" rx="30" fill="#efe2d1" stroke="none"/>
  <text x="600" y="540" font-size="40" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">scripts = fait, references = cadre, assets = fournit</text>
</svg>""",
        "governance": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="135" y="120" width="930" height="150" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="135" y="330" width="930" height="190" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <circle cx="220" cy="195" r="42" fill="#b85724"/>
  <text x="220" y="207" font-size="30" text-anchor="middle" fill="#fffaf5" font-family="Aptos" font-weight="700">1</text>
  <text x="290" y="188" font-size="30" fill="#212427" font-family="Aptos" font-weight="700">Pas d auto-rewrite</text>
  <text x="290" y="225" font-size="22" fill="#5b6268" font-family="Aptos">Sinon perte de stabilite, de tracabilite et de propriete.</text>
  <circle cx="220" cy="425" r="42" fill="#2d5c7a"/>
  <text x="220" y="437" font-size="30" text-anchor="middle" fill="#fffaf5" font-family="Aptos" font-weight="700">2</text>
  <text x="290" y="418" font-size="30" fill="#212427" font-family="Aptos" font-weight="700">Amelioration explicite</text>
  <text x="290" y="455" font-size="22" fill="#5b6268" font-family="Aptos">On observe, on formalise, on corrige le skill, puis on relance.</text>
  <path d="M820 270 L820 330" stroke="#b85724" stroke-width="8" fill="none"/>
  <polygon points="808,320 820,345 832,320" fill="#b85724"/>
</svg>""",
        "compare_tools": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="82" y="112" width="444" height="408" rx="32" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="674" y="112" width="444" height="408" rx="32" fill="#efe2d1" stroke="#dccabf" stroke-width="3"/>
  <circle cx="176" cy="196" r="58" fill="#2d5c7a"/>
  <rect x="148" y="172" width="56" height="48" rx="10" fill="#fffaf5"/>
  <rect x="160" y="186" width="32" height="8" fill="#2d5c7a"/>
  <rect x="160" y="200" width="24" height="8" fill="#2d5c7a"/>
  <text x="304" y="184" font-size="30" fill="#2d5c7a" font-family="Aptos" font-weight="700">Code Companion V3</text>
  <text x="304" y="224" font-size="22" fill="#212427" font-family="Aptos">IDE assistant sécurisé</text>
  <text x="304" y="262" font-size="20" fill="#5b6268" font-family="Aptos">Azure entreprise • PyCharm • VS Code</text>
  <circle cx="768" cy="196" r="58" fill="#b85724"/>
  <path d="M742 215 L792 165 M742 165 L792 215" stroke="#fffaf5" stroke-width="10" stroke-linecap="round"/>
  <text x="896" y="184" font-size="30" fill="#b85724" font-family="Aptos" font-weight="700">Codex CLI</text>
  <text x="896" y="224" font-size="22" fill="#212427" font-family="Aptos">Agent terminal orienté repo</text>
  <text x="896" y="262" font-size="20" fill="#5b6268" font-family="Aptos">skills • AGENTS.md • scripts</text>
  <rect x="126" y="330" width="356" height="126" rx="24" fill="#f8efe5" stroke="none"/>
  <rect x="718" y="330" width="356" height="126" rx="24" fill="#fff7ef" stroke="none"/>
  <text x="304" y="382" font-size="24" text-anchor="middle" fill="#212427" font-family="Aptos">idéal pour modifier vite</text>
  <text x="304" y="420" font-size="24" text-anchor="middle" fill="#212427" font-family="Aptos">et itérer dans l’IDE</text>
  <text x="896" y="382" font-size="24" text-anchor="middle" fill="#212427" font-family="Aptos">idéal pour analyser, patcher,</text>
  <text x="896" y="420" font-size="24" text-anchor="middle" fill="#212427" font-family="Aptos">vérifier et enchaîner</text>
</svg>""",
        "when_to_use": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <path d="M600 110 L1080 560 H120 L600 110 Z" fill="#efe2d1"/>
  <circle cx="300" cy="470" r="76" fill="#2d5c7a"/>
  <text x="300" y="478" font-size="28" text-anchor="middle" fill="#fffaf5" font-family="Aptos" font-weight="700">IDE</text>
  <circle cx="900" cy="470" r="76" fill="#b85724"/>
  <text x="900" y="478" font-size="28" text-anchor="middle" fill="#fffaf5" font-family="Aptos" font-weight="700">CLI</text>
  <text x="600" y="214" font-size="34" text-anchor="middle" fill="#212427" font-family="Aptos" font-weight="700">Choix simple</text>
  <text x="600" y="256" font-size="24" text-anchor="middle" fill="#5b6268" font-family="Aptos">plus local et immédiat à gauche • plus agentique et structuré à droite</text>
  <text x="300" y="600" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">questions locales</text>
  <text x="300" y="632" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">édition rapide</text>
  <text x="900" y="600" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">tâches longues</text>
  <text x="900" y="632" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">workflow complet</text>
</svg>""",
        "multi_provider": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="84" y="122" width="310" height="392" rx="30" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="445" y="122" width="310" height="392" rx="30" fill="#f9f4ed" stroke="#dccabf" stroke-width="3"/>
  <rect x="806" y="122" width="310" height="392" rx="30" fill="#efe2d1" stroke="#dccabf" stroke-width="3"/>
  <text x="239" y="168" font-size="16" text-anchor="middle" fill="#5b6268" font-family="Aptos" font-weight="700">BUILD / GATEWAY</text>
  <text x="600" y="168" font-size="16" text-anchor="middle" fill="#5b6268" font-family="Aptos" font-weight="700">BUY / UI</text>
  <text x="961" y="168" font-size="16" text-anchor="middle" fill="#5b6268" font-family="Aptos" font-weight="700">AGENT</text>
  <text x="239" y="216" font-size="29" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">Portail multi-LLM</text>
  <text x="239" y="265" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">OpenRouter</text>
  <text x="239" y="297" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">Portkey</text>
  <text x="239" y="329" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">LibreChat</text>
  <text x="239" y="382" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">vite à lancer</text>
  <text x="239" y="412" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">gateway pro</text>
  <text x="239" y="442" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">UI interne</text>
  <text x="600" y="216" font-size="29" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">TypingMind vs direct</text>
  <text x="600" y="278" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">TypingMind = couche UI</text>
  <text x="600" y="310" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">OpenAI / Anthropic = moteur</text>
  <text x="600" y="368" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">bon pour usage humain rapide</text>
  <text x="600" y="398" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">pas base produit idéale</text>
  <text x="961" y="216" font-size="29" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">Agent de codage</text>
  <text x="961" y="278" font-size="22" text-anchor="middle" fill="#212427" font-family="Aptos">pas équivalent</text>
  <text x="961" y="338" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">TypingMind : assistant UI</text>
  <text x="961" y="368" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">Codex CLI : repo + fichiers</text>
  <text x="961" y="398" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">tests + scripts + boucle</text>
  <rect x="148" y="548" width="904" height="76" rx="24" fill="#fffaf5" stroke="none"/>
  <text x="600" y="580" font-size="18" text-anchor="middle" fill="#5b6268" font-family="Aptos" font-weight="700">RECOMMANDATION</text>
  <text x="600" y="613" font-size="24" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Pour AGILab : APIs directes + gateway + UI maison</text>
</svg>""",
        "productivity_gain": """<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="680" viewBox="0 0 1200 680">
  <rect width="1200" height="680" fill="#f7f3eb"/>
  <rect x="72" y="132" width="300" height="352" rx="28" fill="#fffaf5" stroke="#dccabf" stroke-width="3"/>
  <rect x="450" y="102" width="300" height="382" rx="28" fill="#f9f4ed" stroke="#dccabf" stroke-width="3"/>
  <rect x="828" y="72" width="300" height="412" rx="28" fill="#efe2d1" stroke="#dccabf" stroke-width="3"/>
  <text x="222" y="174" font-size="18" text-anchor="middle" fill="#5b6268" font-family="Aptos" font-weight="700">NIVEAU 1</text>
  <text x="600" y="144" font-size="18" text-anchor="middle" fill="#5b6268" font-family="Aptos" font-weight="700">NIVEAU 2</text>
  <text x="978" y="114" font-size="18" text-anchor="middle" fill="#5b6268" font-family="Aptos" font-weight="700">NIVEAU 3</text>
  <text x="222" y="214" font-size="31" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Chatbot</text>
  <text x="222" y="250" font-size="21" text-anchor="middle" fill="#212427" font-family="Aptos">propose du code</text>
  <text x="222" y="302" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">copier-coller manuel</text>
  <text x="222" y="332" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">tests a lancer soi-meme</text>
  <text x="222" y="362" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">contexte souvent incomplet</text>
  <text x="222" y="430" font-size="18" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">gain: faible</text>
  <text x="600" y="184" font-size="31" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Plugin IDE</text>
  <text x="600" y="220" font-size="21" text-anchor="middle" fill="#212427" font-family="Aptos">agit dans le fichier courant</text>
  <text x="600" y="272" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">editions plus rapides</text>
  <text x="600" y="302" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">meilleur contexte local</text>
  <text x="600" y="332" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">mais orchestration encore humaine</text>
  <text x="600" y="430" font-size="18" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">gain: moyen</text>
  <text x="978" y="154" font-size="31" text-anchor="middle" fill="#2d5c7a" font-family="Aptos" font-weight="700">Agent multi-skill</text>
  <text x="978" y="190" font-size="21" text-anchor="middle" fill="#212427" font-family="Aptos">enchaîne un vrai workflow</text>
  <text x="978" y="242" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">ouvre le bon contexte repo</text>
  <text x="978" y="272" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">lance scripts et tests</text>
  <text x="978" y="302" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">verifie avant de conclure</text>
  <text x="978" y="332" font-size="20" text-anchor="middle" fill="#5b6268" font-family="Aptos">capitalise les skills du repo</text>
  <text x="978" y="430" font-size="18" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">gain: fort</text>
  <path d="M242 540 C430 500, 770 460, 1018 392" stroke="#b85724" stroke-width="10" fill="none" stroke-linecap="round"/>
  <polygon points="1000,382 1035,386 1011,412" fill="#b85724"/>
  <text x="630" y="580" font-size="30" text-anchor="middle" fill="#b85724" font-family="Aptos" font-weight="700">Le gain de productivite vient de l orchestration retiree a l humain.</text>
</svg>""",
    }

    rendered: dict[str, Path] = {}
    for name, svg in diagrams.items():
        svg_path = DIAGRAM_DIR / f"{name}.svg"
        # Canonical SVGs under docs/source remain the source of truth once they
        # have been manually tuned. The embedded strings only seed missing files.
        if not svg_path.exists():
            _write(svg_path, svg.replace("#dcca bf", "#dccabf"))
        rendered[name] = _render_png(svg_path)
    return rendered


def _set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_band(slide) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def _textbox(slide, left, top, width, height, text="", font_size=20, bold=False, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    return box


def _title(slide, title: str, subtitle: str | None = None) -> None:
    _set_bg(slide)
    _add_band(slide)
    _textbox(slide, Inches(0.7), Inches(0.65), Inches(11.7), Inches(0.75), title, 25, True)
    if subtitle:
        _textbox(slide, Inches(0.72), Inches(1.3), Inches(11.2), Inches(0.45), subtitle, 13, False, MUTED)
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.88), Inches(2.2), Inches(0.07)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def _bullets(slide, left, top, width, height, items: list[str], font_size: int = 21) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK
        p.font.name = "Aptos"
        p.space_after = Pt(14)
        p.bullet = True


def _hero_card(slide, title: str, subtitle: str) -> None:
    _set_bg(slide)
    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(0.6), Inches(12.0), Inches(6.1)
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = PANEL
    hero.line.color.rgb = LINE
    _textbox(slide, Inches(1.0), Inches(1.08), Inches(8.6), Inches(1.5), title, 28, True)
    _textbox(slide, Inches(1.02), Inches(2.5), Inches(8.8), Inches(0.9), subtitle, 20, False, ACCENT_2)
    _textbox(slide, Inches(1.02), Inches(5.85), Inches(5.0), Inches(0.35), "Version plus visuelle", 12, False, MUTED)


def _visual_slide(slide, title: str, subtitle: str, image_path: Path, bullets: list[str], note: str | None = None) -> None:
    _title(slide, title, subtitle)
    slide.shapes.add_picture(str(image_path), Inches(0.85), Inches(2.15), width=Inches(6.0))
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(2.1), Inches(5.3), Inches(3.95)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = LINE
    _bullets(slide, Inches(7.45), Inches(2.55), Inches(4.7), Inches(3.2), bullets, 19)
    if note:
        _textbox(slide, Inches(0.9), Inches(6.25), Inches(11.3), Inches(0.3), note, 14, False, MUTED)


def build_deck(output_path: Path) -> None:
    diagrams = _build_diagrams()
    prs = Presentation()
    prs.slide_width = int(SLIDE_W)
    prs.slide_height = int(SLIDE_H)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _hero_card(
        slide,
        "Pourquoi Codex CLI devient plus performant avec les skills",
        "Lecture par l'exemple AGILab",
    )
    slide.shapes.add_picture(str(diagrams["story"]), Inches(8.15), Inches(1.0), width=Inches(3.7))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _visual_slide(
        slide,
        "1. Point de depart dans AGILab",
        "Un repo vaste, plusieurs facons de lancer et de tester",
        diagrams["story"],
        [
            "UI Streamlit + pages",
            "apps et run configs",
            "plusieurs conventions a respecter",
        ],
        "Dans AGILab, le cout n'est pas seulement d'ecrire du code. Il faut agir au bon endroit.",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _visual_slide(
        slide,
        "2. Comment Codex s'oriente dans AGILab",
        "Il reduit progressivement l'espace de recherche",
        diagrams["flow"],
        [
            "d'abord la description du skill",
            "puis le SKILL.md utile",
            "puis seulement les references et scripts necessaires",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _visual_slide(
        slide,
        "3. Exemple AGILab concret",
        "Corriger une page ou un formulaire sans casser le reste",
        diagrams["agilab_case"],
        [
            "AGENTS.md impose les guardrails",
            "le skill rappelle le workflow",
            "les scripts rendent les etapes repetitives fiables",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _visual_slide(
        slide,
        "4. Les 4 couches utiles a memoriser",
        "Dans AGILab, elles ne jouent pas le meme role",
        diagrams["layers"],
        [
            "AGENTS.md: regles du repo",
            "Skill: procedure",
            "Plugin: distribution",
            "MCP: outils externes",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _visual_slide(
        slide,
        "5. Ce qu'il y a dans un skill",
        "Lecture rapide des dossiers",
        diagrams["folders"],
        [
            "scripts: execute",
            "references: cadrent",
            "assets: fournissent",
        ],
        "Memo simple: scripts = fait, references = explique, assets = fournit.",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "6. Traduit en langage AGILab", "Ce que le skill apporte vraiment")
    _textbox(slide, Inches(1.0), Inches(2.2), Inches(5.6), Inches(3.3), "Quand on travaille sur AGILab,\nle skill evite de re-decouvrir a chaque fois:\n\n- comment lancer\n- quoi verifier\n- quoi ne pas toucher\n- quels tests cibler", 22)
    _bullets(
        slide,
        Inches(7.0),
        Inches(2.4),
        Inches(4.7),
        Inches(2.7),
        [
            "uv partout",
            "logs d'installation d'abord",
            "fixe local avant shared core",
        ],
        21,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _visual_slide(
        slide,
        "7. Pourquoi Codex ne reecrit pas ces skills tout seul",
        "Dans AGILab, ce serait risquer de casser la regle du jeu",
        diagrams["governance"],
        [
            "Eviter la derive silencieuse",
            "Preserver la reproductibilite",
            "Laisser la mise a jour explicite",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "8. Ce que Codex peut faire a la place", "Amelioration explicite, version AGILab")
    _bullets(
        slide,
        Inches(1.0),
        Inches(2.45),
        Inches(10.8),
        Inches(2.5),
        [
            "Observer une friction recurrente",
            "La transformer en skill clair",
            "L'aligner avec AGENTS.md",
            "Le faire evoluer volontairement",
        ],
        23,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title(slide, "9. Conclusion", "L'histoire complete en une phrase")
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(2.15), Inches(11.3), Inches(3.0)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = LINE
    _textbox(slide, Inches(1.35), Inches(2.8), Inches(10.2), Inches(0.9), "Dans AGILab, les skills rendent Codex plus utile parce qu'ils lui apprennent comment bien travailler ici.", 24, True, INK)
    _textbox(slide, Inches(1.35), Inches(4.0), Inches(10.2), Inches(0.7), "Et ils restent des artefacts humains, donc mis a jour volontairement plutot qu'automatiquement.", 18, False, MUTED)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir) / output_path.name
        prs.save(str(tmp_path))
        output_path.write_bytes(tmp_path.read_bytes())


def main() -> int:
    build_deck(OUTPUT)
    print(OUTPUT.resolve())
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
