#!/usr/bin/env python3
from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FLIGHT_DECK_NAME = "AGILAB_Flight_First_Proof_Slides.pptx"
UAV_DECK_NAME = "AGILAB_UAV_Full_Tour_Slides.pptx"
ARTIFACT_DIR = ROOT / "artifacts" / "demo_media"
DOCS_DIR = ROOT / "docs" / "source"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(11, 19, 29)
BG_ALT = RGBColor(18, 32, 45)
SURFACE = RGBColor(247, 243, 238)
SURFACE_ALT = RGBColor(227, 233, 240)
SURFACE_DARK = RGBColor(28, 45, 61)
TEXT_LIGHT = RGBColor(244, 240, 235)
TEXT_MUTED = RGBColor(171, 183, 196)
TEXT_DARK = RGBColor(25, 33, 42)
TEXT_SOFT_DARK = RGBColor(87, 96, 109)
ACCENT = RGBColor(232, 110, 58)
ACCENT_2 = RGBColor(84, 185, 214)
ACCENT_3 = RGBColor(245, 191, 91)
LINE_LIGHT = RGBColor(85, 102, 121)
LINE_DARK = RGBColor(214, 204, 193)

FONT_UI = "Aptos"
FONT_DISPLAY = "Aptos Display"


def add_full_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_run_font(
    run,
    *,
    size: int,
    color: RGBColor,
    bold: bool = False,
    name: str = FONT_UI,
) -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_decor(slide) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0),
        Inches(0.0),
        Inches(13.333),
        Inches(0.16),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    corner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(11.6),
        Inches(-0.8),
        Inches(2.8),
        Inches(2.8),
    )
    corner.fill.solid()
    corner.fill.fore_color.rgb = BG_ALT
    corner.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(10.9),
        Inches(5.6),
        Inches(2.3),
        Inches(2.3),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = SURFACE_DARK
    accent.line.fill.background()


def add_kicker(
    slide,
    text: str,
    *,
    left,
    top,
    width,
    color: RGBColor = ACCENT_2,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.25))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text.upper()
    _set_run_font(run, size=10, color=color, bold=True)
    p.space_after = 0
    p.space_before = 0


def add_title_block(
    slide,
    title: str,
    subtitle: str,
    *,
    left,
    top,
    width,
    title_size: int = 26,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _set_run_font(r, size=title_size, color=TEXT_LIGHT, bold=True, name=FONT_DISPLAY)
    p.space_after = Pt(3)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = subtitle
    _set_run_font(r2, size=13, color=TEXT_MUTED)
    p2.space_before = Pt(2)


def add_panel(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill: RGBColor = SURFACE,
    line: RGBColor = LINE_DARK,
    rounded: bool = True,
):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if rounded
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def add_shadow_panel(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill: RGBColor = SURFACE,
    line: RGBColor = LINE_DARK,
) -> None:
    shadow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + Inches(0.07),
        top + Inches(0.08),
        width,
        height,
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = BG_ALT
    shadow.line.fill.background()
    add_panel(slide, left, top, width, height, fill=fill, line=line)


def add_tag(
    slide,
    text: str,
    *,
    left,
    top,
    width=None,
    fill: RGBColor = ACCENT,
    text_color: RGBColor = TEXT_LIGHT,
) -> None:
    if width is None:
        width = Inches(max(1.2, min(2.8, 0.26 + len(text) * 0.09)))
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        Inches(0.34),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_run_font(r, size=9, color=text_color, bold=True)


def add_story_card(
    slide,
    number: str,
    title: str,
    body: str,
    *,
    left,
    top,
    width,
    height,
) -> None:
    add_shadow_panel(slide, left, top, width, height)
    add_tag(slide, number, left=left + Inches(0.22), top=top + Inches(0.18), width=Inches(0.48), fill=ACCENT_2)

    tb = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.58), width - Inches(0.44), height - Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, size=15, color=TEXT_DARK, bold=True, name=FONT_DISPLAY)
    p.space_after = Pt(3)

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    _set_run_font(r2, size=11, color=TEXT_SOFT_DARK)


def add_statement_band(
    slide,
    text: str,
    *,
    left,
    top,
    width,
    height=Inches(0.86),
    fill: RGBColor = SURFACE_DARK,
) -> None:
    band = add_panel(slide, left, top, width, height, fill=fill, line=fill, rounded=True)
    band.line.fill.background()
    tb = slide.shapes.add_textbox(left + Inches(0.24), top + Inches(0.18), width - Inches(0.48), height - Inches(0.24))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    _set_run_font(r, size=12, color=TEXT_LIGHT, bold=False)


def add_paragraph_list(
    slide,
    items: list[str],
    *,
    left,
    top,
    width,
    height,
    font_size: int = 14,
    color: RGBColor = TEXT_DARK,
    lead_color: RGBColor | None = None,
) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if lead_color is not None:
            r0 = p.add_run()
            r0.text = "— "
            _set_run_font(r0, size=font_size, color=lead_color, bold=True)
        r = p.add_run()
        r.text = item
        _set_run_font(r, size=font_size, color=color)
        p.space_after = Pt(8)


def add_feature_card(
    slide,
    title: str,
    body: str,
    *,
    left,
    top,
    width,
    height,
    accent: RGBColor = ACCENT,
) -> None:
    add_shadow_panel(slide, left, top, width, height, fill=SURFACE, line=LINE_DARK)
    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        Inches(0.13),
        height,
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    tb = slide.shapes.add_textbox(left + Inches(0.28), top + Inches(0.2), width - Inches(0.5), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run_font(r, size=14, color=TEXT_DARK, bold=True, name=FONT_DISPLAY)
    p.space_after = Pt(2)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    _set_run_font(r2, size=11, color=TEXT_SOFT_DARK)


def _image_cover_bytes(
    image_path: Path,
    *,
    width_emu: int,
    height_emu: int,
    darken: float = 1.0,
) -> BytesIO:
    with Image.open(image_path) as image:
        image = image.convert("RGB")
        source_ratio = image.width / image.height
        target_ratio = width_emu / height_emu

        if source_ratio > target_ratio:
            crop_width = int(round(image.height * target_ratio))
            offset_x = max(0, (image.width - crop_width) // 2)
            box = (offset_x, 0, offset_x + crop_width, image.height)
        else:
            crop_height = int(round(image.width / target_ratio))
            offset_y = max(0, (image.height - crop_height) // 2)
            box = (0, offset_y, image.width, offset_y + crop_height)

        cropped = image.crop(box)
        if darken != 1.0:
            cropped = ImageEnhance.Brightness(cropped).enhance(darken)

        output = BytesIO()
        cropped.save(output, format="PNG")
        output.seek(0)
        return output


def add_image_card(
    slide,
    image_path: Path,
    *,
    left,
    top,
    width,
    height,
    label: str,
    caption: str | None = None,
    accent: RGBColor = ACCENT,
    darken: float = 1.0,
) -> None:
    shadow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left + Inches(0.07),
        top + Inches(0.08),
        width,
        height,
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = BG_ALT
    shadow.line.fill.background()

    frame = add_panel(slide, left, top, width, height, fill=SURFACE, line=LINE_DARK, rounded=True)
    frame.fill.fore_color.rgb = SURFACE
    inner_left = left + Inches(0.12)
    inner_top = top + Inches(0.12)
    inner_width = width - Inches(0.24)
    inner_height = height - Inches(0.24)

    if image_path.exists():
        slide.shapes.add_picture(
            _image_cover_bytes(
                image_path,
                width_emu=int(inner_width),
                height_emu=int(inner_height),
                darken=darken,
            ),
            inner_left,
            inner_top,
            width=inner_width,
            height=inner_height,
        )
    else:
        placeholder = add_panel(
            slide,
            inner_left,
            inner_top,
            inner_width,
            inner_height,
            fill=SURFACE_ALT,
            line=LINE_DARK,
            rounded=False,
        )
        placeholder.line.width = Pt(0.8)
        add_paragraph_list(
            slide,
            [f"Missing image: {image_path.name}"],
            left=inner_left + Inches(0.2),
            top=inner_top + Inches(0.22),
            width=inner_width - Inches(0.4),
            height=Inches(0.5),
            font_size=11,
            color=TEXT_SOFT_DARK,
        )

    add_tag(slide, label, left=left + Inches(0.2), top=top - Inches(0.16), fill=accent, text_color=TEXT_LIGHT)
    if caption:
        caption_box = slide.shapes.add_textbox(left + Inches(0.18), top + height - Inches(0.42), width - Inches(0.36), Inches(0.22))
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = caption
        _set_run_font(r, size=10, color=TEXT_LIGHT, bold=True)


def add_step_pills(
    slide,
    steps: list[tuple[str, RGBColor]],
    *,
    left,
    top,
) -> None:
    x = left
    for label, color in steps:
        width = Inches(max(1.25, min(2.1, 0.5 + len(label) * 0.09)))
        add_tag(slide, label, left=x, top=top, width=width, fill=color, text_color=TEXT_LIGHT)
        x += width + Inches(0.12)


def _canonical_docs_dir() -> Path | None:
    canonical = ROOT.parent / "thales_agilab" / "docs" / "source"
    return canonical if canonical.parent.exists() else None


def _default_output_targets(filename: str, artifact_name: str) -> list[Path]:
    targets = [ARTIFACT_DIR / artifact_name, DOCS_DIR / filename]
    canonical = _canonical_docs_dir()
    if canonical is not None:
        targets.append(canonical / filename)
    return targets


def _write_presentation(prs: Presentation, output_targets: list[Path]) -> list[Path]:
    buffer = BytesIO()
    prs.save(buffer)
    payload = buffer.getvalue()
    written: list[Path] = []
    for output in output_targets:
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_bytes(payload)
        written.append(output)
    return written


def build_flight_deck(*, output_targets: list[Path] | None = None) -> list[Path]:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page_shots = ROOT / "docs" / "source" / "_static" / "page-shots"
    core_overview = page_shots / "core-pages-overview.png"
    project_page = page_shots / "project-page.png"
    orchestrate_page = page_shots / "orchestrate-page.png"
    analysis_page = page_shots / "analysis-page.png"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_decor(slide)
    add_kicker(slide, "newcomer-safe proof", left=Inches(0.85), top=Inches(0.72), width=Inches(4.5))
    add_title_block(
        slide,
        "AGILAB First Proof",
        "Slideshow companion for the newcomer-safe `flight_project` story",
        left=Inches(0.85),
        top=Inches(0.95),
        width=Inches(5.2),
        title_size=28,
    )
    add_story_card(
        slide,
        "1",
        "One app",
        "Stay on `flight_project` and keep the first story narrow.",
        left=Inches(0.85),
        top=Inches(2.45),
        width=Inches(1.85),
        height=Inches(1.6),
    )
    add_story_card(
        slide,
        "2",
        "One path",
        "Use PROJECT, then ORCHESTRATE, without side-shell glue.",
        left=Inches(2.92),
        top=Inches(2.45),
        width=Inches(1.85),
        height=Inches(1.6),
    )
    add_story_card(
        slide,
        "3",
        "One result",
        "End on visible evidence in ANALYSIS rather than raw logs.",
        left=Inches(4.99),
        top=Inches(2.45),
        width=Inches(1.85),
        height=Inches(1.6),
    )
    add_statement_band(
        slide,
        "Goal: prove the safest AGILAB path once before branching into notebooks, clusters, or packaged installs.",
        left=Inches(0.85),
        top=Inches(4.45),
        width=Inches(5.95),
        fill=SURFACE_DARK,
    )
    add_image_card(
        slide,
        project_page,
        left=Inches(7.55),
        top=Inches(1.1),
        width=Inches(5.0),
        height=Inches(2.65),
        label="PROJECT",
        caption="Select the built-in app and keep its context visible.",
        accent=ACCENT,
    )
    add_image_card(
        slide,
        orchestrate_page,
        left=Inches(7.15),
        top=Inches(2.9),
        width=Inches(4.9),
        height=Inches(2.65),
        label="ORCHESTRATE",
        caption="Run the local proof through one controlled path.",
        accent=ACCENT_2,
    )
    add_image_card(
        slide,
        analysis_page,
        left=Inches(8.15),
        top=Inches(4.25),
        width=Inches(4.2),
        height=Inches(2.15),
        label="ANALYSIS",
        caption="Visible result, not infrastructure noise.",
        accent=ACCENT_3,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_decor(slide)
    add_kicker(slide, "story logic", left=Inches(0.85), top=Inches(0.72), width=Inches(2.5))
    add_title_block(
        slide,
        "Why This Path",
        "The first proof is deliberately smaller than the full tour.",
        left=Inches(0.85),
        top=Inches(0.95),
        width=Inches(5.2),
    )
    add_feature_card(
        slide,
        "Visible app context",
        "PROJECT is not just a launcher. It keeps the selected app legible while you run it.",
        left=Inches(0.85),
        top=Inches(2.1),
        width=Inches(4.15),
        height=Inches(1.1),
        accent=ACCENT,
    )
    add_feature_card(
        slide,
        "Controlled execution",
        "ORCHESTRATE packages and runs the local path without asking the viewer to infer shell state.",
        left=Inches(0.85),
        top=Inches(3.38),
        width=Inches(4.15),
        height=Inches(1.1),
        accent=ACCENT_2,
    )
    add_feature_card(
        slide,
        "Readable proof",
        "ANALYSIS is the payoff: a result a technical audience can parse in seconds.",
        left=Inches(0.85),
        top=Inches(4.66),
        width=Inches(4.15),
        height=Inches(1.1),
        accent=ACCENT_3,
    )
    add_image_card(
        slide,
        core_overview,
        left=Inches(5.45),
        top=Inches(1.55),
        width=Inches(6.9),
        height=Inches(4.5),
        label="UI FLOW",
        caption="PROJECT → ORCHESTRATE → ANALYSIS is the newcomer-safe route through the UI.",
        accent=ACCENT_2,
    )
    add_statement_band(
        slide,
        "Narration: this is the smallest AGILAB story that still proves real workflow control and a visible result.",
        left=Inches(0.85),
        top=Inches(6.18),
        width=Inches(11.55),
        fill=SURFACE_DARK,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_decor(slide)
    add_kicker(slide, "execution path", left=Inches(0.85), top=Inches(0.72), width=Inches(2.5))
    add_title_block(
        slide,
        "Select And Run",
        "PROJECT sets the app context; ORCHESTRATE executes the local proof.",
        left=Inches(0.85),
        top=Inches(0.95),
        width=Inches(7.0),
    )
    add_step_pills(
        slide,
        [("1 PROJECT", ACCENT), ("2 ORCHESTRATE", ACCENT_2)],
        left=Inches(0.85),
        top=Inches(1.9),
    )
    add_image_card(
        slide,
        project_page,
        left=Inches(0.85),
        top=Inches(2.35),
        width=Inches(5.55),
        height=Inches(3.3),
        label="PROJECT",
        caption="Choose `flight_project` and keep the scenario visible.",
        accent=ACCENT,
    )
    add_image_card(
        slide,
        orchestrate_page,
        left=Inches(6.95),
        top=Inches(2.35),
        width=Inches(5.55),
        height=Inches(3.3),
        label="ORCHESTRATE",
        caption="INSTALL, then EXECUTE through the local path.",
        accent=ACCENT_2,
    )
    add_panel(slide, Inches(0.85), Inches(6.0), Inches(11.65), Inches(0.9), fill=SURFACE, line=LINE_DARK)
    add_paragraph_list(
        slide,
        [
            "Use `flight_project`: the first proof is not about breadth; it is about getting one clean end-to-end run.",
            "Keep the visual story simple and operator-visible so the evidence reads immediately.",
        ],
        left=Inches(1.1),
        top=Inches(6.18),
        width=Inches(11.1),
        height=Inches(0.5),
        font_size=12,
        color=TEXT_DARK,
        lead_color=ACCENT,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_decor(slide)
    add_kicker(slide, "proof state", left=Inches(0.85), top=Inches(0.72), width=Inches(2.5))
    add_title_block(
        slide,
        "End On Evidence",
        "ANALYSIS is where the first proof becomes concrete.",
        left=Inches(0.85),
        top=Inches(0.95),
        width=Inches(5.0),
    )
    add_feature_card(
        slide,
        "Proof condition",
        "Fresh output under `~/log/execute/flight/` plus a readable ANALYSIS result means the newcomer hurdle is cleared.",
        left=Inches(0.85),
        top=Inches(2.05),
        width=Inches(4.25),
        height=Inches(1.3),
        accent=ACCENT,
    )
    add_feature_card(
        slide,
        "Why this matters",
        "The point is not that AGILAB ran somewhere. The point is that it ended on visible evidence a reviewer can inspect.",
        left=Inches(0.85),
        top=Inches(3.58),
        width=Inches(4.25),
        height=Inches(1.3),
        accent=ACCENT_2,
    )
    add_tag(slide, "VISIBLE RESULT", left=Inches(0.85), top=Inches(5.25), width=Inches(1.5), fill=ACCENT_3, text_color=TEXT_DARK)
    add_statement_band(
        slide,
        "Close the story on the result screen. That is what makes the first proof credible in one glance.",
        left=Inches(0.85),
        top=Inches(5.68),
        width=Inches(4.6),
        fill=SURFACE_DARK,
    )
    add_image_card(
        slide,
        analysis_page,
        left=Inches(5.75),
        top=Inches(1.45),
        width=Inches(6.55),
        height=Inches(4.95),
        label="ANALYSIS",
        caption="Readable result screen, centered as the payoff of the first proof.",
        accent=ACCENT_3,
    )

    outputs = output_targets or _default_output_targets(
        FLIGHT_DECK_NAME,
        "agilab_flight_first_proof_slides.pptx",
    )
    return _write_presentation(prs, outputs)


def build_uav_deck(*, output_targets: list[Path] | None = None) -> list[Path]:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page_shots = ROOT / "docs" / "source" / "_static" / "page-shots"
    core_overview = page_shots / "core-pages-overview.png"
    project_page = page_shots / "project-page.png"
    orchestrate_page = page_shots / "orchestrate-page.png"
    pipeline_page = page_shots / "pipeline-page.png"
    analysis_page = page_shots / "analysis-page.png"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_decor(slide)
    add_kicker(slide, "public workflow story", left=Inches(0.85), top=Inches(0.72), width=Inches(4.0))
    add_title_block(
        slide,
        "AGILAB Full Tour",
        "Slideshow companion for the `uav_relay_queue_project` public workflow story.",
        left=Inches(0.85),
        top=Inches(0.95),
        width=Inches(5.3),
        title_size=28,
    )
    add_step_pills(
        slide,
        [
            ("PROJECT", ACCENT),
            ("ORCHESTRATE", ACCENT_2),
            ("PIPELINE", ACCENT_3),
            ("ANALYSIS", RGBColor(152, 126, 224)),
        ],
        left=Inches(0.85),
        top=Inches(2.15),
    )
    add_statement_band(
        slide,
        "Goal: show the four-page AGILAB story without asking the viewer to replay a hosted runtime or watch a video.",
        left=Inches(0.85),
        top=Inches(2.78),
        width=Inches(5.7),
        fill=SURFACE_DARK,
    )
    add_story_card(
        slide,
        "A",
        "Replayable path",
        "The run is not just executable once; the workflow is made explicit.",
        left=Inches(0.85),
        top=Inches(4.0),
        width=Inches(2.65),
        height=Inches(1.55),
    )
    add_story_card(
        slide,
        "B",
        "Readable payoff",
        "End on queue or route evidence that a technical reviewer can scan fast.",
        left=Inches(3.82),
        top=Inches(4.0),
        width=Inches(2.65),
        height=Inches(1.55),
    )
    add_image_card(
        slide,
        project_page,
        left=Inches(7.2),
        top=Inches(1.18),
        width=Inches(4.9),
        height=Inches(2.45),
        label="PROJECT",
        caption="App and scenario are visible from the start.",
        accent=ACCENT,
    )
    add_image_card(
        slide,
        pipeline_page,
        left=Inches(7.75),
        top=Inches(2.95),
        width=Inches(4.9),
        height=Inches(2.45),
        label="PIPELINE",
        caption="The run becomes an explicit tracked workflow.",
        accent=ACCENT_3,
    )
    add_image_card(
        slide,
        analysis_page,
        left=Inches(8.4),
        top=Inches(4.55),
        width=Inches(4.2),
        height=Inches(1.75),
        label="ANALYSIS",
        caption="Readable end-state evidence.",
        accent=ACCENT_2,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_decor(slide)
    add_kicker(slide, "workflow contract", left=Inches(0.85), top=Inches(0.72), width=Inches(3.2))
    add_title_block(
        slide,
        "The Full Workflow Story",
        "This is the stronger README and public-tour narrative.",
        left=Inches(0.85),
        top=Inches(0.95),
        width=Inches(5.3),
    )
    add_image_card(
        slide,
        core_overview,
        left=Inches(0.85),
        top=Inches(1.95),
        width=Inches(6.45),
        height=Inches(4.2),
        label="FOUR PAGE FLOW",
        caption="The story is the page sequence, not isolated screenshots.",
        accent=ACCENT_2,
    )
    add_feature_card(
        slide,
        "Explicit context",
        "PROJECT anchors the app and scenario instead of hiding setup inside an implicit shell state.",
        left=Inches(7.65),
        top=Inches(2.0),
        width=Inches(4.5),
        height=Inches(1.1),
        accent=ACCENT,
    )
    add_feature_card(
        slide,
        "Replayable workflow",
        "PIPELINE shows that execution steps and artifacts can be inspected and rerun, not only re-described.",
        left=Inches(7.65),
        top=Inches(3.28),
        width=Inches(4.5),
        height=Inches(1.1),
        accent=ACCENT_3,
    )
    add_feature_card(
        slide,
        "Technical payoff",
        "ANALYSIS lands on queue or route evidence that justifies the workflow story.",
        left=Inches(7.65),
        top=Inches(4.56),
        width=Inches(4.5),
        height=Inches(1.1),
        accent=ACCENT_2,
    )
    add_statement_band(
        slide,
        "Narration: AGILAB turns a lightweight experiment into an explicit, reviewable workflow instead of a one-off script path.",
        left=Inches(0.85),
        top=Inches(6.28),
        width=Inches(11.3),
        fill=SURFACE_DARK,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_decor(slide)
    add_kicker(slide, "page sequence", left=Inches(0.85), top=Inches(0.72), width=Inches(2.8))
    add_title_block(
        slide,
        "Project To Pipeline",
        "PROJECT and ORCHESTRATE establish the run; PIPELINE makes it explicit.",
        left=Inches(0.85),
        top=Inches(0.95),
        width=Inches(7.1),
    )
    add_image_card(
        slide,
        project_page,
        left=Inches(0.85),
        top=Inches(2.05),
        width=Inches(3.8),
        height=Inches(3.0),
        label="1 PROJECT",
        caption="Select `uav_relay_queue_project` and keep the scenario visible.",
        accent=ACCENT,
    )
    add_image_card(
        slide,
        orchestrate_page,
        left=Inches(4.78),
        top=Inches(2.05),
        width=Inches(3.8),
        height=Inches(3.0),
        label="2 ORCHESTRATE",
        caption="Trigger the run through one packaged path.",
        accent=ACCENT_2,
    )
    add_image_card(
        slide,
        pipeline_page,
        left=Inches(8.71),
        top=Inches(2.05),
        width=Inches(3.8),
        height=Inches(3.0),
        label="3 PIPELINE",
        caption="Show that the execution path is tracked and replayable.",
        accent=ACCENT_3,
    )
    add_panel(slide, Inches(0.85), Inches(5.45), Inches(11.65), Inches(1.0), fill=SURFACE, line=LINE_DARK)
    add_paragraph_list(
        slide,
        [
            "The deck works only if the viewer reads these as one route through the product, not as disconnected screens.",
            "That is why PIPELINE matters in the full tour and not in the newcomer-safe first proof.",
        ],
        left=Inches(1.08),
        top=Inches(5.67),
        width=Inches(11.1),
        height=Inches(0.55),
        font_size=12,
        color=TEXT_DARK,
        lead_color=ACCENT,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_bg(slide)
    add_decor(slide)
    add_kicker(slide, "technical payoff", left=Inches(0.85), top=Inches(0.72), width=Inches(2.8))
    add_title_block(
        slide,
        "End On Queue Evidence",
        "ANALYSIS carries the technical payoff of the full tour.",
        left=Inches(0.85),
        top=Inches(0.95),
        width=Inches(5.4),
    )
    add_image_card(
        slide,
        analysis_page,
        left=Inches(0.85),
        top=Inches(1.85),
        width=Inches(7.15),
        height=Inches(4.7),
        label="ANALYSIS",
        caption="Queue or route evidence is where the public tour becomes credible.",
        accent=ACCENT_2,
    )
    add_feature_card(
        slide,
        "What the viewer should retain",
        "AGILAB did not just execute a scenario. It kept the workflow explicit all the way to a readable technical result.",
        left=Inches(8.35),
        top=Inches(2.0),
        width=Inches(3.8),
        height=Inches(1.55),
        accent=ACCENT,
    )
    add_feature_card(
        slide,
        "Why this deck exists",
        "It gives the README and docs a skimmable public demo surface without requiring a hosted runtime or video playback.",
        left=Inches(8.35),
        top=Inches(3.78),
        width=Inches(3.8),
        height=Inches(1.55),
        accent=ACCENT_3,
    )
    add_tag(slide, "FULL TOUR", left=Inches(8.35), top=Inches(5.72), width=Inches(1.2), fill=ACCENT_2, text_color=TEXT_LIGHT)
    add_statement_band(
        slide,
        "The point of the full tour is not breadth for its own sake. It is explicit workflow plus visible queue evidence.",
        left=Inches(8.35),
        top=Inches(6.1),
        width=Inches(3.8),
        height=Inches(0.58),
        fill=SURFACE_DARK,
    )

    outputs = output_targets or _default_output_targets(
        UAV_DECK_NAME,
        "agilab_uav_full_tour_slides.pptx",
    )
    return _write_presentation(prs, outputs)


def build() -> dict[str, list[Path]]:
    return {
        "flight": build_flight_deck(),
        "uav": build_uav_deck(),
    }


def main() -> int:
    outputs = build()
    for deck_outputs in outputs.values():
        for output in deck_outputs:
            print(output.resolve())
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
