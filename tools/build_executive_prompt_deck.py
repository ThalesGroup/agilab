#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs/source/TP Prompt Ingénierie Executive.pptx"
ARTIFACT = ROOT / "artifacts/tp_prompt_ingenierie_executive.pptx"
DIAG = ROOT / "docs/source/diagrams/codex_skills"

BG = RGBColor(250, 247, 242)
PAPER = RGBColor(255, 252, 248)
INK = RGBColor(28, 31, 35)
MUTED = RGBColor(95, 101, 107)
ACCENT = RGBColor(191, 90, 36)
ACCENT_2 = RGBColor(33, 84, 114)
ACCENT_3 = RGBColor(223, 214, 204)
DARK = RGBColor(24, 34, 44)
WHITE = RGBColor(255, 255, 255)


def set_bg(slide, dark: bool = False) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK if dark else BG


def panel(slide, left, top, width, height, *, fill=PAPER, line=ACCENT_3):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def pill(slide, left, top, width, height, text, *, fill=RGBColor(245, 236, 225), color=ACCENT, size=12):
    shape = panel(slide, left, top, width, height, fill=fill, line=fill)
    textbox(slide, left, top + Inches(0.03), width, height, text, size=size, color=color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def textbox(slide, left, top, width, height, text="", *, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Aptos"
    return box


def title_block(slide, title: str, subtitle: str):
    textbox(slide, Inches(0.82), Inches(0.72), Inches(11.0), Inches(0.55), title, size=27, color=ACCENT_2, bold=True)
    textbox(slide, Inches(0.84), Inches(1.35), Inches(11.0), Inches(0.34), subtitle, size=13, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(1.82), Inches(1.8), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def slide_prompt(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, "1. Le vrai rôle du Prompt Engineering", "Il cadre un coding harness : modèle, agent loop et runtime supports.")
    textbox(slide, Inches(4.7), Inches(2.02), Inches(3.9), Inches(0.36), "Coding harness", size=22, color=DARK, bold=True, align=PP_ALIGN.CENTER)

    left_panel = panel(slide, Inches(0.74), Inches(2.42), Inches(2.72), Inches(3.1), fill=WHITE, line=DARK)
    center_panel = panel(slide, Inches(4.02), Inches(2.42), Inches(4.82), Inches(3.1), fill=WHITE, line=DARK)
    right_panel = panel(slide, Inches(9.4), Inches(2.42), Inches(3.18), Inches(3.1), fill=WHITE, line=DARK)
    for shp in (left_panel, center_panel, right_panel):
        shp.line.width = Pt(2.0)

    textbox(slide, Inches(1.0), Inches(2.72), Inches(2.16), Inches(0.28), "Model family", size=21, color=DARK, align=PP_ALIGN.CENTER)
    panel(slide, Inches(1.14), Inches(3.34), Inches(1.92), Inches(0.88), fill=RGBColor(250, 217, 176), line=RGBColor(250, 217, 176))
    textbox(slide, Inches(1.3), Inches(3.64), Inches(1.6), Inches(0.2), "Base LLM", size=19, color=DARK, align=PP_ALIGN.CENTER)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(1.95), Inches(4.28), Inches(0.28), Inches(0.28))
    shape.rotation = 90
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    panel(slide, Inches(1.14), Inches(4.62), Inches(1.92), Inches(0.88), fill=RGBColor(250, 217, 176), line=RGBColor(250, 217, 176))
    textbox(slide, Inches(1.3), Inches(4.84), Inches(1.6), Inches(0.36), "Reasoning\nmodel", size=18, color=DARK, align=PP_ALIGN.CENTER)

    textbox(slide, Inches(5.16), Inches(2.72), Inches(2.52), Inches(0.28), "Agent loop", size=21, color=DARK, align=PP_ALIGN.CENTER)
    loop_boxes = [
        (Inches(4.46), Inches(3.66), "Inspect"),
        (Inches(6.58), Inches(3.66), "Choose"),
        (Inches(4.46), Inches(4.76), "Observe"),
        (Inches(6.58), Inches(4.76), "Act"),
    ]
    for left, top, label in loop_boxes:
        panel(slide, left, top, Inches(1.55), Inches(0.56), fill=RGBColor(197, 225, 245), line=RGBColor(197, 225, 245))
        textbox(slide, left, top + Inches(0.15), Inches(1.55), Inches(0.18), label, size=18, color=DARK, align=PP_ALIGN.CENTER)
    arrows = [
        (Inches(6.05), Inches(3.86), Inches(0.24), Inches(0.16), 0),
        (Inches(7.13), Inches(4.26), Inches(0.16), Inches(0.24), 90),
        (Inches(6.05), Inches(5.0), Inches(0.24), Inches(0.16), 180),
        (Inches(5.13), Inches(4.22), Inches(0.16), Inches(0.24), 270),
    ]
    for left, top, width, height, rotation in arrows:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
        shape.rotation = rotation
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK
        shape.line.fill.background()

    textbox(slide, Inches(9.7), Inches(2.72), Inches(2.56), Inches(0.28), "Runtime supports", size=20, color=DARK, align=PP_ALIGN.CENTER)
    runtime = [
        (Inches(9.68), Inches(3.56), "Repo context"),
        (Inches(11.1), Inches(3.56), "Tools"),
        (Inches(9.68), Inches(4.34), "Permissions"),
        (Inches(11.1), Inches(4.34), "Memory"),
        (Inches(9.68), Inches(5.12), "Cache"),
        (Inches(11.1), Inches(5.12), "Execution"),
    ]
    for left, top, label in runtime:
        panel(slide, left, top, Inches(1.16), Inches(0.48), fill=RGBColor(226, 208, 246), line=RGBColor(226, 208, 246))
        textbox(slide, left, top + Inches(0.12), Inches(1.16), Inches(0.18), label, size=14, color=DARK, align=PP_ALIGN.CENTER)

    for left in [Inches(3.56), Inches(9.0)]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, Inches(4.02), Inches(0.3), Inches(0.24))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK
        shape.line.fill.background()


def slide_landscape(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, "2. Le landscape agentique en 4 couches", "Le modèle calcule. Le framework organise. L’agent agit. L’UI expose.")
    for i, (title, body, fill) in enumerate([
        ("MODÈLE", "génère\ntexte ou code", RGBColor(234, 243, 252)),
        ("FRAMEWORK", "structure état\net outils", RGBColor(240, 236, 251)),
        ("AGENT", "choisit,\nenchaîne,\nvérifie", RGBColor(249, 238, 228)),
        ("UI / IDE", "rend l’action\nvisible", RGBColor(252, 243, 234)),
    ]):
        left = Inches(0.86 + i * 3.12)
        panel(slide, left, Inches(2.35), Inches(2.7), Inches(2.45), fill=fill)
        pill(slide, left + Inches(0.52), Inches(2.63), Inches(1.62), Inches(0.34), title, fill=WHITE, color=ACCENT_2, size=10)
        textbox(slide, left + Inches(0.2), Inches(3.35), Inches(2.3), Inches(0.85), body, size=17, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    for left in [Inches(3.25), Inches(6.37), Inches(9.49)]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, Inches(3.45), Inches(0.45), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
    panel(slide, Inches(1.18), Inches(5.12), Inches(11.0), Inches(0.75), fill=WHITE)
    textbox(slide, Inches(1.45), Inches(5.36), Inches(10.45), Inches(0.22),
            "Dans AGILab, AGENTS.md + skills + scripts renforcent surtout la couche agent.", size=18, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)


def slide_agilab(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, "3. Pourquoi AGILab est un bon cas d’école", "Dans un repo riche, le coût d’orchestration devient visible tout de suite.")
    for left, title, body in [
        (Inches(1.0), "Demande", "corriger une page\nsans casser le repo"),
        (Inches(5.05), "Guidage", "AGENTS.md\n+ skills"),
        (Inches(9.1), "Sortie", "patch fiable\net validation ciblée"),
    ]:
        panel(slide, left, Inches(2.5), Inches(2.9), Inches(2.15), fill=WHITE)
        textbox(slide, left + Inches(0.22), Inches(2.86), Inches(2.45), Inches(0.28), title, size=20, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, left + Inches(0.22), Inches(3.48), Inches(2.45), Inches(0.62), body, size=17, color=INK, align=PP_ALIGN.CENTER)
    for left in [Inches(4.25), Inches(8.3)]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, Inches(3.45), Inches(0.45), Inches(0.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
    panel(slide, Inches(1.05), Inches(5.25), Inches(11.1), Inches(0.72), fill=WHITE, line=ACCENT_3)
    textbox(slide, Inches(1.35), Inches(5.5), Inches(10.5), Inches(0.22),
            "Même workload, benchmark explicite : PandasWorker = process ; PolarsWorker = threads.", size=17, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)


def slide_codex(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, "4. Pourquoi Codex CLI devient performant", "Son intérêt : agir avec le bon contexte, pas tout relire ni tout improviser.")
    steps = ["métadonnées", "SKILL.md", "scripts utiles", "exécution"]
    for i, step in enumerate(steps):
        left = Inches(0.95 + i * 2.95)
        panel(slide, left, Inches(2.65), Inches(2.2), Inches(1.15), fill=WHITE)
        pill(slide, left + Inches(0.8), Inches(2.3), Inches(0.58), Inches(0.34), str(i + 1), fill=ACCENT, color=WHITE)
        textbox(slide, left + Inches(0.15), Inches(3.08), Inches(1.9), Inches(0.22), step, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(3.0 + i * 2.95), Inches(3.08), Inches(0.36), Inches(0.24))
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT
            shape.line.fill.background()
    panel(slide, Inches(1.2), Inches(5.1), Inches(10.9), Inches(0.72), fill=RGBColor(244, 236, 224), line=RGBColor(244, 236, 224))
    textbox(slide, Inches(1.45), Inches(5.35), Inches(10.3), Inches(0.22),
            "Il n’ouvre que ce qui est utile, au moment utile.", size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


def slide_skills(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, "5. Ce que les skills apportent vraiment", "Ils rendent le workflow stable, réutilisable et plus fiable.")
    for i, (title, tag) in enumerate([("scripts/", "exécute"), ("references/", "cadre"), ("assets/", "fournit")]):
        left = Inches(1.0 + i * 4.05)
        panel(slide, left, Inches(2.35), Inches(3.3), Inches(2.4), fill=WHITE if i != 2 else RGBColor(249, 243, 236))
        textbox(slide, left + Inches(0.2), Inches(2.72), Inches(2.9), Inches(0.26), title, size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        pill(slide, left + Inches(0.82), Inches(3.2), Inches(1.65), Inches(0.38), tag)


def slide_governance(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, "6. Pourquoi ils ne sont pas mis à jour automatiquement", "Plus d’autonomie ne veut pas dire réécriture silencieuse du repo.")
    panel(slide, Inches(1.1), Inches(2.45), Inches(4.6), Inches(2.5), fill=WHITE)
    panel(slide, Inches(7.6), Inches(2.45), Inches(4.6), Inches(2.5), fill=RGBColor(249, 243, 236))
    pill(slide, Inches(2.45), Inches(2.78), Inches(1.85), Inches(0.4), "AUTONOMIE")
    pill(slide, Inches(8.45), Inches(2.78), Inches(2.85), Inches(0.4), "REÉCRITURE NON")
    textbox(slide, Inches(1.4), Inches(3.45), Inches(4.0), Inches(0.6), "oui pour choisir,\nlancer et vérifier", size=20, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, Inches(7.9), Inches(3.45), Inches(4.0), Inches(0.6), "non pour modifier seul\nles règles du repo", size=20, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)


def slide_compare(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    title_block(slide, "7. Code Companion V3 vs Codex CLI", "Le plugin fait gagner du temps. L’agent fait gagner des étapes de travail.")
    for left, title, body, fill, color in [
        (Inches(0.95), "Code Companion V3", "édition locale rapide\nbon pour itérer", WHITE, ACCENT_2),
        (Inches(6.82), "Codex CLI", "scripts + vérifications\nenchaîne le workflow", RGBColor(249, 243, 236), ACCENT),
    ]:
        panel(slide, left, Inches(2.35), Inches(5.0), Inches(2.55), fill=fill)
        textbox(slide, left + Inches(0.25), Inches(2.78), Inches(4.5), Inches(0.28), title, size=23, color=color, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, left + Inches(0.25), Inches(3.52), Inches(4.5), Inches(0.6), body, size=19, color=INK, align=PP_ALIGN.CENTER)


def slide_takeaways(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    textbox(slide, Inches(0.82), Inches(0.8), Inches(11.0), Inches(0.55), "8. À retenir", size=27, color=ACCENT_2, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(1.82), Inches(1.8), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    panel(slide, Inches(1.0), Inches(2.0), Inches(11.2), Inches(1.35), fill=WHITE, line=ACCENT_3)
    textbox(slide, Inches(1.3), Inches(2.15), Inches(10.6), Inches(0.75),
            "Le vrai gain vient de l’orchestration retirée à l’humain.", size=24, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    cards = [
        ("Prompt", "cadre l’action"),
        ("Plugin", "accélère localement"),
        ("Agent", "enchaîne le workflow"),
    ]
    for i, (title, body) in enumerate(cards):
        left = 1.1 + i * 3.75
        panel(slide, Inches(left), Inches(4.0), Inches(3.2), Inches(1.35), fill=WHITE, line=ACCENT_3)
        textbox(slide, Inches(left + 0.18), Inches(4.22), Inches(2.84), Inches(0.22), title, size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, Inches(left + 0.18), Inches(4.66), Inches(2.84), Inches(0.3), body, size=17, color=INK, align=PP_ALIGN.CENTER)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_prompt(prs)
    slide_landscape(prs)
    slide_agilab(prs)
    slide_codex(prs)
    slide_skills(prs)
    slide_governance(prs)
    slide_compare(prs)
    slide_takeaways(prs)
    return prs


def _first_text(slide) -> str:
    for shape in slide.shapes:
        text = getattr(shape, "text", "").strip()
        if text and text != "RADIO ACADEMIE":
            return text
    return ""


def validate_outline(prs: Presentation) -> None:
    expected = [
        "1. Le vrai rôle du Prompt Engineering",
        "2. Le landscape agentique en 4 couches",
        "3. Pourquoi AGILab est un bon cas d’école",
        "4. Pourquoi Codex CLI devient performant",
        "5. Ce que les skills apportent vraiment",
        "6. Pourquoi ils ne sont pas mis à jour automatiquement",
        "7. Code Companion V3 vs Codex CLI",
        "8. À retenir",
    ]
    actual = [_first_text(slide) for slide in prs.slides]
    if actual != expected:
        raise ValueError(f"Unexpected executive outline: {actual}")


def main() -> int:
    prs = build()
    validate_outline(prs)
    ARTIFACT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(ARTIFACT))
    prs.save(str(OUT))
    print(ARTIFACT)
    print("slides", len(prs.slides))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
