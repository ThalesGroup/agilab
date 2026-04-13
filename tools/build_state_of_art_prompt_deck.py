#!/usr/bin/env python3
from __future__ import annotations

import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs/source/TP Prompt Ingénierie.pptx"
DIAG = ROOT / "docs/source/diagrams/codex_skills"

BG = RGBColor(250, 247, 242)
PAPER = RGBColor(255, 252, 248)
INK = RGBColor(26, 31, 36)
MUTED = RGBColor(94, 101, 107)
ACCENT = RGBColor(191, 90, 36)
ACCENT_2 = RGBColor(33, 84, 114)
ACCENT_3 = RGBColor(223, 214, 204)
DARK = RGBColor(23, 34, 45)
WHITE = RGBColor(255, 255, 255)


def bg(slide, dark: bool = False) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK if dark else BG


def box(slide, left, top, width, height, text="", *, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Aptos"
    return tb


def bullets(slide, left, top, width, height, items, *, size=19, color=INK):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Aptos"
        p.space_after = Pt(12)
        p.bullet = True
    return tb


def panel(slide, left, top, width, height, *, fill=PAPER, line=ACCENT_3):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def pill(slide, left, top, width, height, text, *, fill=RGBColor(245, 236, 225), color=ACCENT, size=12):
    shape = panel(slide, left, top, width, height, fill=fill, line=fill)
    box(slide, left, top + Inches(0.03), width, height, text, size=size, color=color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def section(slide, title: str, subtitle: str) -> None:
    bg(slide)
    box(slide, Inches(0.82), Inches(0.66), Inches(11.5), Inches(0.5), title, size=27, color=ACCENT_2, bold=True)
    box(slide, Inches(0.84), Inches(1.28), Inches(11.0), Inches(0.3), subtitle, size=13, color=MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(1.76), Inches(1.8), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def slide_problem(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(
        slide,
        "1. Le vrai rôle du Prompt Engineering",
        "Le prompt ne cadre pas seulement le modèle. Il cadre le système d'exécution : modèle, boucle agentique et runtime.",
    )
    box(
        slide,
        Inches(4.18),
        Inches(2.0),
        Inches(5.05),
        Inches(0.52),
        "Système cadré par le prompt",
        size=20,
        color=DARK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    left_panel = panel(slide, Inches(0.7), Inches(2.48), Inches(2.8), Inches(3.28), fill=WHITE, line=DARK)
    center_panel = panel(slide, Inches(4.08), Inches(2.48), Inches(4.72), Inches(3.28), fill=WHITE, line=DARK)
    right_panel = panel(slide, Inches(9.38), Inches(2.48), Inches(3.25), Inches(3.28), fill=WHITE, line=DARK)
    for shp in (left_panel, center_panel, right_panel):
        shp.line.width = Pt(2.2)

    box(slide, Inches(0.97), Inches(2.78), Inches(2.25), Inches(0.32), "Couche modèle", size=22, color=DARK, bold=False, align=PP_ALIGN.CENTER)
    panel(slide, Inches(1.12), Inches(3.46), Inches(1.98), Inches(0.94), fill=RGBColor(250, 217, 176), line=RGBColor(250, 217, 176))
    box(slide, Inches(1.22), Inches(3.78), Inches(1.78), Inches(0.24), "LLM de base", size=20, color=DARK, bold=False, align=PP_ALIGN.CENTER)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(1.93), Inches(4.46), Inches(0.34), Inches(0.34))
    shape.rotation = 90
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    panel(slide, Inches(1.12), Inches(4.88), Inches(1.98), Inches(0.94), fill=RGBColor(250, 217, 176), line=RGBColor(250, 217, 176))
    box(slide, Inches(1.22), Inches(5.12), Inches(1.78), Inches(0.42), "LLM\nreasoning", size=20, color=DARK, bold=False, align=PP_ALIGN.CENTER)

    box(slide, Inches(5.1), Inches(2.78), Inches(2.72), Inches(0.32), "Boucle agentique", size=22, color=DARK, bold=False, align=PP_ALIGN.CENTER)
    loop_boxes = [
        (Inches(4.56), Inches(3.84), "Inspecte"),
        (Inches(6.58), Inches(3.84), "Choisit"),
        (Inches(4.56), Inches(5.02), "Observe"),
        (Inches(6.58), Inches(5.02), "Agit"),
    ]
    for left, top, label in loop_boxes:
        panel(slide, left, top, Inches(1.48), Inches(0.62), fill=RGBColor(197, 225, 245), line=RGBColor(197, 225, 245))
        box(slide, left, top + Inches(0.17), Inches(1.48), Inches(0.2), label, size=19, color=DARK, bold=False, align=PP_ALIGN.CENTER)
    arrows = [
        (Inches(6.06), Inches(4.06), Inches(0.28), Inches(0.18), 0),
        (Inches(7.16), Inches(4.48), Inches(0.18), Inches(0.28), 90),
        (Inches(6.06), Inches(5.28), Inches(0.28), Inches(0.18), 180),
        (Inches(5.16), Inches(4.46), Inches(0.18), Inches(0.28), 270),
    ]
    for left, top, width, height, rotation in arrows:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, width, height)
        shape.rotation = rotation
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK
        shape.line.fill.background()

    box(slide, Inches(9.5), Inches(2.78), Inches(2.98), Inches(0.32), "Cadre d'exécution", size=20, color=DARK, bold=False, align=PP_ALIGN.CENTER)
    runtime = [
        (Inches(9.6), Inches(3.78), "Contexte\nrepo"),
        (Inches(11.06), Inches(3.78), "Outils"),
        (Inches(9.6), Inches(4.62), "Permissions"),
        (Inches(11.06), Inches(4.62), "Mémoire"),
        (Inches(9.6), Inches(5.46), "Cache"),
        (Inches(11.06), Inches(5.46), "Exécution"),
    ]
    for left, top, label in runtime:
        panel(slide, left, top, Inches(1.38), Inches(0.58), fill=RGBColor(226, 208, 246), line=RGBColor(226, 208, 246))
        box(slide, left, top + Inches(0.09), Inches(1.38), Inches(0.34), label, size=14, color=DARK, bold=False, align=PP_ALIGN.CENTER)

    for left in [Inches(3.56), Inches(8.94)]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, Inches(4.12), Inches(0.34), Inches(0.28))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK
        shape.line.fill.background()

    box(
        slide,
        Inches(1.0),
        Inches(6.1),
        Inches(11.3),
        Inches(0.42),
        "Le Prompt Engineering sert à cadrer la boucle, les outils, les règles d'exécution et le bon niveau de contexte.",
        size=16,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_landscape(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "2. Le landscape agentique en 4 couches", "Le modèle calcule. Le framework organise. L’agent agit. L’UI expose.")
    layers = [
        ("MODÈLE", "produit texte,\ncode, plan,\nrésumé", RGBColor(234, 243, 252)),
        ("FRAMEWORK", "structure état,\noutils, mémoire,\nrouting", RGBColor(240, 236, 251)),
        ("AGENT", "choisit,\nenchaîne,\nvérifie", RGBColor(249, 238, 228)),
        ("UI / IDE", "rend le système\nvisible et\npilotable", RGBColor(252, 243, 234)),
    ]
    for i, (title, body, fill) in enumerate(layers):
        left = Inches(0.86 + i * 3.12)
        panel(slide, left, Inches(2.18), Inches(2.7), Inches(2.75), fill=fill)
        pill(slide, left + Inches(0.54), Inches(2.46), Inches(1.62), Inches(0.34), title, fill=WHITE, color=ACCENT_2, size=10)
        box(slide, left + Inches(0.2), Inches(3.06), Inches(2.3), Inches(1.02), body, size=19, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    for left in [Inches(3.3), Inches(6.42), Inches(9.54)]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, Inches(3.28), Inches(0.48), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
    panel(slide, Inches(0.98), Inches(5.0), Inches(5.25), Inches(0.92), fill=WHITE)
    box(slide, Inches(1.18), Inches(5.18), Inches(4.86), Inches(0.2), "Message clé", size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.18), Inches(5.44), Inches(4.86), Inches(0.32),
        "Le gain visible n’est pas dans la couche modèle seule.", size=19, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
    panel(slide, Inches(6.55), Inches(5.0), Inches(5.78), Inches(0.92), fill=RGBColor(249, 243, 236))
    box(slide, Inches(6.78), Inches(5.18), Inches(5.32), Inches(0.2), "Cas AGILab", size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(6.78), Inches(5.44), Inches(5.32), Inches(0.32),
        "AGENTS.md + skills + scripts renforcent surtout la couche agent.", size=19, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)


def slide_ladder(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "3. La montée réelle de productivité", "Même modèle. Trois niveaux de valeur très différents.")
    stairs = [
        (Inches(0.78), Inches(3.08), Inches(3.35), Inches(1.7), RGBColor(244, 240, 235)),
        (Inches(4.05), Inches(2.46), Inches(3.6), Inches(2.32), RGBColor(242, 235, 227)),
        (Inches(7.48), Inches(1.78), Inches(4.12), Inches(3.0), RGBColor(239, 227, 214)),
    ]
    for left, top, width, height, fill in stairs:
        panel(slide, left, top, width, height, fill=fill, line=fill)
    cards = [
        (Inches(0.95), Inches(2.26), Inches(2.95), Inches(2.52), "NIVEAU 1", "Chatbot", "répond", ["copier-coller", "tests manuels"], "GAIN LIMITÉ", WHITE, 22, 15, 14, Inches(1.44), Inches(0.40), Inches(0.38)),
        (Inches(4.24), Inches(1.86), Inches(3.16), Inches(2.82), "NIVEAU 2", "Plugin IDE", "édite", ["contexte local", "itération rapide"], "GAIN MOYEN", RGBColor(255, 250, 244), 22, 15, 14, Inches(1.48), Inches(0.44), Inches(0.40)),
        (Inches(7.82), Inches(1.42), Inches(3.72), Inches(3.22), "NIVEAU 3", "Agent\nmulti-skill", "orchestre", ["repo + skills", "scripts + vérification"], "GAIN FORT", RGBColor(255, 248, 238), 20, 15, 14, Inches(1.50), Inches(0.44), Inches(0.40)),
    ]
    for left, top, width, height, level, title, verb, items, gain, fill, title_size, verb_size, body_size, body_top, body_height, gain_offset in cards:
        panel(slide, left, top, width, height, fill=fill)
        pill(slide, left + Inches(0.24), top + Inches(0.18), Inches(1.18), Inches(0.34), level, size=11)
        box(slide, left + Inches(0.24), top + Inches(0.58), width - Inches(0.48), Inches(0.56), title, size=title_size, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
        box(slide, left + Inches(0.24), top + Inches(1.14), width - Inches(0.48), Inches(0.24), verb, size=verb_size, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        box(slide, left + Inches(0.28), top + body_top, width - Inches(0.56), body_height, "\n".join(f"• {item}" for item in items), size=body_size, color=MUTED, align=PP_ALIGN.CENTER)
        pill(slide, left + (width - Inches(1.58)) / 2, top + height - gain_offset, Inches(1.58), Inches(0.32), gain, fill=RGBColor(245, 236, 225), size=10)
    ramp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(1.5), Inches(5.1), Inches(9.9), Inches(0.42))
    ramp.fill.solid()
    ramp.fill.fore_color.rgb = ACCENT
    ramp.line.fill.background()
    box(slide, Inches(1.55), Inches(5.54), Inches(9.8), Inches(0.48),
        "Plus on retire l’orchestration à l’humain, plus le gain change d’échelle.", size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


def slide_agilab(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "4. Pourquoi AGILab est un bon cas d’école", "Dans un repo riche, le coût d’orchestration devient visible immédiatement.")
    stages = [
        (Inches(0.95), "Demande", "corriger une page\nou un formulaire"),
        (Inches(4.55), "Guidage", "prompt + AGENTS.md\n+ skills du repo"),
        (Inches(8.15), "Sortie", "patch plus fiable\net validation ciblée"),
    ]
    for left, title, body in stages:
        panel(slide, left, Inches(2.35), Inches(2.8), Inches(2.25), fill=WHITE)
        box(slide, left + Inches(0.2), Inches(2.72), Inches(2.4), Inches(0.3), title, size=22, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
        box(slide, left + Inches(0.18), Inches(3.35), Inches(2.45), Inches(0.72), body, size=18, color=INK, align=PP_ALIGN.CENTER)
    for left in [Inches(3.85), Inches(7.45)]:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, Inches(3.23), Inches(0.42), Inches(0.38))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
    panel(slide, Inches(1.0), Inches(5.05), Inches(11.0), Inches(0.9), fill=RGBColor(245, 236, 225), line=RGBColor(245, 236, 225))
    box(slide, Inches(1.25), Inches(5.24), Inches(10.5), Inches(0.46),
        "AGILab montre bien où un agent apporte de la valeur : il rappelle comment bien travailler dans ce repo.", size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    panel(slide, Inches(1.0), Inches(6.08), Inches(11.0), Inches(0.7), fill=WHITE, line=ACCENT_3)
    box(slide, Inches(1.3), Inches(6.24), Inches(10.4), Inches(0.26),
        "Benchmark lisible : PandasWorker = process, PolarsWorker = threads, même workload, exécution rendue explicite.", size=15, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)


def slide_codex(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "5. Pourquoi Codex CLI devient performant", "Son intérêt : agir avec le bon contexte, pas tout lire ni tout improviser.")
    steps = ["métadonnées repo", "SKILL.md", "scripts utiles", "exécution"]
    for i, step in enumerate(steps):
        left = Inches(0.88 + i * 3.0)
        panel(slide, left, Inches(2.55), Inches(2.42), Inches(1.34), fill=WHITE)
        pill(slide, left + Inches(0.96), Inches(2.14), Inches(0.62), Inches(0.36), str(i + 1), fill=ACCENT, color=WHITE)
        box(slide, left + Inches(0.16), Inches(3.03), Inches(2.22), Inches(0.34), step, size=17, color=INK, bold=True, align=PP_ALIGN.CENTER)
        if i < 3:
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(3.47 + i * 3.0), Inches(3.08), Inches(0.34), Inches(0.24))
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT
            shape.line.fill.background()
    panel(slide, Inches(1.2), Inches(5.0), Inches(10.8), Inches(0.86), fill=RGBColor(244, 236, 224), line=RGBColor(244, 236, 224))
    box(slide, Inches(1.45), Inches(5.27), Inches(10.2), Inches(0.24),
        "Progressive disclosure : Codex n’ouvre que ce qui est utile, au moment utile.", size=19, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)


def slide_skills(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "6. Ce qu’un skill apporte vraiment", "Le skill transforme une tâche vague en workflow reproductible.")
    cols = [
        ("scripts/", "exécute", "CLI déterministe\nrun_eval.py\ncollect_metrics.py"),
        ("references/", "cadre", "protocoles\nchecklists\nmétriques"),
        ("assets/", "fournit", "templates\nexemples d’input\nressources locales"),
    ]
    for i, (title, tag, body) in enumerate(cols):
        left = Inches(0.92 + i * 4.03)
        panel(slide, left, Inches(2.28), Inches(3.35), Inches(2.85), fill=WHITE if i != 2 else RGBColor(249, 243, 236))
        box(slide, left + Inches(0.2), Inches(2.62), Inches(2.95), Inches(0.3), title, size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        pill(slide, left + Inches(0.85), Inches(3.12), Inches(1.6), Inches(0.4), tag)
        box(slide, left + Inches(0.3), Inches(3.78), Inches(2.75), Inches(1.2), body, size=18, color=MUTED, align=PP_ALIGN.CENTER)


def slide_governance(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "7. Pourquoi Codex ne met pas les skills à jour tout seul", "Plus d’autonomie ne veut pas dire réécriture silencieuse du repo.")
    panel(slide, Inches(1.0), Inches(2.25), Inches(5.0), Inches(2.7), fill=WHITE)
    pill(slide, Inches(1.3), Inches(2.55), Inches(1.7), Inches(0.42), "AUTONOMIE")
    box(slide, Inches(1.35), Inches(3.15), Inches(4.3), Inches(0.9),
        "oui pour lire,\nchoisir, lancer et vérifier", size=21, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
    panel(slide, Inches(7.25), Inches(2.25), Inches(5.0), Inches(2.7), fill=RGBColor(249, 243, 236))
    pill(slide, Inches(7.58), Inches(2.55), Inches(2.35), Inches(0.42), "RÉÉCRITURE SILENCIEUSE")
    box(slide, Inches(7.55), Inches(3.15), Inches(4.35), Inches(0.9),
        "non pour modifier seul\nles règles du repo", size=21, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.2), Inches(5.42), Inches(11.0), Inches(0.28),
        "Pourquoi : pas de dérive silencieuse, pas de surprise côté repo, amélioration volontaire.", size=18, color=MUTED, align=PP_ALIGN.CENTER)


def slide_multi_provider(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "8. Portail chatbot multi-provider", "Le bon choix dépend de la couche produit recherchée : UI, gateway, ou agent.")
    slide.shapes.add_picture(str(DIAG / "multi_provider.png"), Inches(0.78), Inches(2.0), width=Inches(11.8))


def slide_compare(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "8. Code Companion V3 vs Codex CLI", "Le plugin fait gagner du temps dans l’IDE. L’agent fait gagner des étapes de travail.")
    cards = [
        (Inches(0.92), "Code Companion V3", "assistant IDE sécurisé", ["édition locale rapide", "bon pour itérer", "l’humain orchestre encore"], WHITE, ACCENT_2),
        (Inches(6.75), "Codex CLI", "agent terminal orienté repo", ["analyse plusieurs fichiers", "lance scripts et vérifications", "enchaîne le workflow"], RGBColor(249, 243, 236), ACCENT),
    ]
    for left, title, subtitle, items, fill, title_color in cards:
        panel(slide, left, Inches(2.1), Inches(5.15), Inches(3.35), fill=fill)
        box(slide, left + Inches(0.28), Inches(2.52), Inches(4.55), Inches(0.34), title, size=24, color=title_color, bold=True)
        box(slide, left + Inches(0.28), Inches(2.95), Inches(4.3), Inches(0.24), subtitle, size=17, color=INK)
        box(slide, left + Inches(0.35), Inches(3.55), Inches(4.2), Inches(1.2), "\n".join(f"• {item}" for item in items), size=18, color=MUTED)
    panel(slide, Inches(2.0), Inches(5.35), Inches(9.35), Inches(0.72), fill=RGBColor(245, 236, 225), line=RGBColor(245, 236, 225))
    box(slide, Inches(2.2), Inches(5.6), Inches(8.95), Inches(0.24),
        "Le plugin fait gagner du temps. L’agent fait gagner des étapes de travail.", size=20, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)


def slide_when(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section(slide, "9. Quand utiliser quoi", "Même équipe, mêmes modèles, mais pas le même niveau d’orchestration.")
    panel(slide, Inches(1.0), Inches(2.35), Inches(4.6), Inches(2.9), fill=WHITE)
    panel(slide, Inches(7.75), Inches(2.35), Inches(4.6), Inches(2.9), fill=RGBColor(249, 243, 236))
    box(slide, Inches(1.28), Inches(2.78), Inches(4.0), Inches(0.35), "Code Companion V3", size=24, color=ACCENT_2, bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(1.35), Inches(3.45), Inches(3.9), Inches(1.0),
        "question locale\nédition rapide\nitération dans l’IDE", size=20, color=INK, align=PP_ALIGN.CENTER)
    box(slide, Inches(8.03), Inches(2.78), Inches(4.0), Inches(0.35), "Codex CLI", size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    box(slide, Inches(8.1), Inches(3.45), Inches(3.9), Inches(1.0),
        "workflow complet\ntâche longue\nvalidation et enchaînement", size=20, color=INK, align=PP_ALIGN.CENTER)
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(5.9), Inches(3.35), Inches(1.45), Inches(0.78))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def slide_takeaways(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide)
    box(slide, Inches(0.82), Inches(0.76), Inches(11.0), Inches(0.55), "10. À retenir", size=27, color=ACCENT_2, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.82), Inches(1.78), Inches(1.8), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    panel(slide, Inches(0.95), Inches(2.0), Inches(11.35), Inches(1.45), fill=WHITE, line=ACCENT_3)
    box(slide, Inches(1.28), Inches(2.12), Inches(10.7), Inches(0.72),
        "Le vrai levier de productivité n’est pas le modèle seul.\nC’est l’orchestration du travail.", size=24, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    cards = [
        ("Prompt Engineering", "cadre modèle, contexte\net sortie attendue"),
        ("Plugin IDE", "accélère l’édition\net l’itération locale"),
        ("Agent multi-skill", "enchaîne workflow,\nscripts et vérification"),
    ]
    for i, (title, text) in enumerate(cards):
        left = 0.95 + i * 3.8
        panel(slide, Inches(left), Inches(4.02), Inches(3.45), Inches(1.7), fill=WHITE, line=ACCENT_3)
        box(slide, Inches(left + 0.2), Inches(4.28), Inches(3.0), Inches(0.28), title, size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        box(slide, Inches(left + 0.2), Inches(4.78), Inches(3.0), Inches(0.6), text, size=17, color=INK, align=PP_ALIGN.CENTER)


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_problem(prs)
    slide_landscape(prs)
    slide_ladder(prs)
    slide_agilab(prs)
    slide_codex(prs)
    slide_skills(prs)
    slide_governance(prs)
    slide_compare(prs)
    slide_when(prs)
    slide_takeaways(prs)
    return prs


def _first_text(slide) -> str:
    for shape in slide.shapes:
        text = getattr(shape, "text", "").strip()
        if text and text != "RADIO ACADEMIE":
            return text
    return ""


def validate_outline(prs: Presentation) -> None:
    expected = [
        "1. Le vrai rôle du Prompt Engineering",
        "2. Le landscape agentique en 4 couches",
        "3. La montée réelle de productivité",
        "4. Pourquoi AGILab est un bon cas d’école",
        "5. Pourquoi Codex CLI devient performant",
        "6. Ce qu’un skill apporte vraiment",
        "7. Pourquoi Codex ne met pas les skills à jour tout seul",
        "8. Code Companion V3 vs Codex CLI",
        "9. Quand utiliser quoi",
        "10. À retenir",
    ]
    actual = [_first_text(slide) for slide in prs.slides]
    if actual != expected:
        raise ValueError(f"Unexpected premium outline: {actual}")


def main() -> int:
    prs = build()
    validate_outline(prs)
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir) / OUT.name
        prs.save(str(tmp))
        OUT.write_bytes(tmp.read_bytes())
    print(OUT)
    print("slides", len(prs.slides))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
