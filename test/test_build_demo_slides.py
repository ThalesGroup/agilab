from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

from pptx import Presentation


MODULE_PATH = Path(__file__).resolve().parents[1] / "tools" / "build_demo_slides.py"
MODULE_SPEC = importlib.util.spec_from_file_location("tools.build_demo_slides", MODULE_PATH)
assert MODULE_SPEC is not None and MODULE_SPEC.loader is not None
demo_slides = importlib.util.module_from_spec(MODULE_SPEC)
sys.modules[MODULE_SPEC.name] = demo_slides
MODULE_SPEC.loader.exec_module(demo_slides)


def _slide_texts(deck_path: Path) -> list[str]:
    prs = Presentation(deck_path)
    assert len(prs.slides) == 4
    return ["\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")) for slide in prs.slides]


def test_build_demo_slides_writes_two_public_decks(tmp_path: Path) -> None:
    flight = tmp_path / demo_slides.FLIGHT_DECK_NAME
    uav = tmp_path / demo_slides.UAV_DECK_NAME

    written_flight = demo_slides.build_flight_deck(output_targets=[flight])
    written_uav = demo_slides.build_uav_deck(output_targets=[uav])

    assert written_flight == [flight]
    assert written_uav == [uav]
    assert flight.exists()
    assert uav.exists()

    flight_texts = _slide_texts(flight)
    assert any("AGILAB First Proof" in text for text in flight_texts)
    assert any("Select And Run" in text for text in flight_texts)
    assert any("End On Evidence" in text for text in flight_texts)

    uav_texts = _slide_texts(uav)
    assert any("AGILAB Full Tour" in text for text in uav_texts)
    assert any("Project To Pipeline" in text for text in uav_texts)
    assert any("End On Queue Evidence" in text for text in uav_texts)
